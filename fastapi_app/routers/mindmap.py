from __future__ import annotations

import asyncio
import json
import re
import shutil
import time
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

import fitz  # PyMuPDF
from fastapi import APIRouter, Body, Depends, File, Form, HTTPException, Request, UploadFile

from dataflow_agent.agentroles import create_agent
from dataflow_agent.logger import get_logger
from dataflow_agent.state import MainRequest, MainState
from fastapi_app.dependencies import AuthUser, get_optional_user, is_auth_configured
from fastapi_app.config.pricing import estimate_mindmap_points
from fastapi_app.services.billing_service import BillingService
from fastapi_app.services.managed_api_service import resolve_llm_credentials
from fastapi_app.utils import _to_outputs_url, get_outputs_root, resolve_outputs_path

router = APIRouter(prefix="/mindmap", tags=["mindmap"])
log = get_logger(__name__)

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

_ALLOWED_SUFFIXES = {
    ".pdf",
    ".doc",
    ".docx",
    ".ppt",
    ".pptx",
    ".txt",
    ".md",
}


def _owner_slug(user: Optional[AuthUser]) -> str:
    raw = ""
    if user is not None:
        raw = (user.email or user.id or "").strip()
    raw = raw or "local"
    return re.sub(r"[^a-zA-Z0-9@._-]+", "_", raw)


def _require_or_resolve_user(user: Optional[AuthUser]) -> Optional[AuthUser]:
    if is_auth_configured() and user is None:
        raise HTTPException(status_code=401, detail="Authentication required")
    return user


def _mindmap_root_for_user(user: Optional[AuthUser]) -> Path:
    owner = _owner_slug(user)
    return (get_outputs_root() / "mindmap" / owner).resolve()


def _resolve_owned_mindmap_path(path_or_url: str, user: Optional[AuthUser]) -> Path:
    resolved = resolve_outputs_path(path_or_url, must_exist=False)
    root = _mindmap_root_for_user(user)
    try:
        resolved.relative_to(root)
    except ValueError as exc:
        raise HTTPException(status_code=403, detail="Mindmap path does not belong to the current user") from exc
    return resolved


async def _parse_one_file(file_path: str) -> Dict[str, str]:
    file_path_obj = Path(file_path)
    filename = file_path_obj.name
    suffix = file_path_obj.suffix.lower()

    if not file_path_obj.exists():
        return {"filename": filename, "content": f"[Error: File not found {file_path}]"}

    raw_content = ""

    try:
        if suffix == ".pdf":
            doc = fitz.open(file_path)
            parts = [page.get_text() for page in doc]
            raw_content = "\n".join(parts)
        elif suffix in {".doc", ".docx"}:
            if Document is None:
                raw_content = "[Error: python-docx not installed]"
            else:
                doc = Document(file_path)
                raw_content = "\n".join(p.text for p in doc.paragraphs)
        elif suffix in {".ppt", ".pptx"}:
            if Presentation is None:
                raw_content = "[Error: python-pptx not installed]"
            else:
                prs = Presentation(file_path)
                slides: List[str] = []
                for index, slide in enumerate(prs.slides, start=1):
                    slide_text = [f"--- Slide {index} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    slides.append("\n".join(slide_text))
                raw_content = "\n".join(slides)
        else:
            raw_content = file_path_obj.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        raw_content = f"[Parse Error: {exc}]"

    return {
        "filename": filename,
        "content": raw_content[:60000] if len(raw_content) > 60000 else raw_content,
    }


async def _parse_files(paths: List[str]) -> List[Dict[str, str]]:
    if not paths:
        return []
    return await asyncio.gather(*[_parse_one_file(path) for path in paths])


def _extract_agent_json_result(state: MainState, role_name: str) -> Any:
    try:
        result = state.agent_results.get(role_name, {}).get("results", {})
        if isinstance(result, dict):
            if "json" in result and result["json"] not in (None, ""):
                return result["json"]
            if "data" in result and result["data"] not in (None, ""):
                return result["data"]
            if "text" in result and isinstance(result["text"], (dict, list)):
                return result["text"]
            if "raw" in result and isinstance(result["raw"], (dict, list)):
                return result["raw"]
            for value in result.values():
                if isinstance(value, (dict, list)):
                    return value
        if isinstance(result, (dict, list)):
            return result
    except Exception:
        return None
    return None


def _coerce_tree_node(node: Any, *, prefix: str, depth: int, max_depth: int) -> Optional[Dict[str, Any]]:
    if depth > max_depth:
        return None

    if isinstance(node, str):
        label = node.strip()
        if not label:
            return None
        return {"id": prefix, "label": label[:80], "summary": "", "children": []}

    if not isinstance(node, dict):
        return None

    label = str(node.get("label") or node.get("title") or node.get("name") or "").strip()
    if not label:
        return None

    summary = str(node.get("summary") or node.get("note") or node.get("description") or "").strip()
    raw_children = node.get("children") or node.get("items") or node.get("nodes") or []
    if not isinstance(raw_children, list):
        raw_children = []

    children: List[Dict[str, Any]] = []
    for index, child in enumerate(raw_children[:8], start=1):
        coerced = _coerce_tree_node(child, prefix=f"{prefix}_{index}", depth=depth + 1, max_depth=max_depth)
        if coerced:
            children.append(coerced)

    return {
        "id": prefix,
        "label": label[:80],
        "summary": summary[:220],
        "children": children,
    }


def _normalize_tree_payload(payload: Any, *, max_depth: int) -> Dict[str, Any]:
    root_candidate: Any = payload
    highlights: List[str] = []

    if isinstance(payload, dict):
        if isinstance(payload.get("highlights"), list):
            highlights = [str(item).strip() for item in payload["highlights"][:6] if str(item).strip()]
        root_candidate = payload.get("root") or payload.get("tree") or payload.get("mindmap") or payload

    root = _coerce_tree_node(root_candidate, prefix="root", depth=0, max_depth=max_depth)
    if root is None:
        root = {
            "id": "root",
            "label": "Mind Map",
            "summary": "No structured output returned by the model.",
            "children": [],
        }

    return {"root": root, "highlights": highlights}


def count_tree_nodes(node: Dict[str, Any]) -> int:
    children = node.get("children") if isinstance(node, dict) else []
    if not isinstance(children, list):
        children = []
    return 1 + sum(count_tree_nodes(child) for child in children if isinstance(child, dict))


def get_tree_depth(node: Dict[str, Any]) -> int:
    children = node.get("children") if isinstance(node, dict) else []
    if not isinstance(children, list) or not children:
        return 1
    depths = [get_tree_depth(child) for child in children if isinstance(child, dict)]
    return 1 + max(depths) if depths else 1


async def _generate_tree(
    *,
    text_blocks: List[Dict[str, str]],
    chat_api_url: str,
    api_key: str,
    model: str,
    max_depth: int,
    language: str,
    style: str,
) -> Dict[str, Any]:
    content_sections = []
    for item in text_blocks:
        filename = item.get("filename", "document")
        content = item.get("content", "").strip()
        if not content:
            continue
        content_sections.append(f"## {filename}\n{content}")

    joined_content = "\n\n".join(content_sections).strip()
    if not joined_content:
        return {
            "root": {
                "id": "root",
                "label": "Mind Map",
                "summary": "No readable content extracted from the input.",
                "children": [],
            },
            "highlights": [],
        }

    prompt = f"""You are an information architect building a mind map tree from source documents.

Return strict JSON only. No markdown fences. No explanation.

Schema:
{{
  "root": {{
    "label": "central topic",
    "summary": "one short sentence",
    "children": [
      {{
        "label": "branch",
        "summary": "optional",
        "children": []
      }}
    ]
  }},
  "highlights": ["short key takeaway", "short key takeaway"]
}}

Rules:
- Language: {language}
- Mind map style hint: {style}
- Max depth: {max_depth}
- Each label must be concise.
- Prefer 3-6 first-level branches.
- Each branch should reflect a meaningful topic, workflow stage, method block, result group, or conclusion.
- Keep summaries short.
- Do not include ids, markdown, HTML, Mermaid, XML, or comments.
- If the material is long, compress aggressively but keep the hierarchy coherent.

Source content:
{joined_content}
"""

    agent = create_agent(
        name="kb_prompt_agent",
        model_name=model,
        chat_api_url=chat_api_url,
        temperature=0.2,
        parser_type="json",
    )

    req = MainRequest(
        language=language,
        chat_api_url=chat_api_url,
        api_key=api_key,
        chat_api_key=api_key,
        model=model,
    )
    temp_state = MainState(request=req)
    res_state = await agent.execute(temp_state, prompt=prompt)
    payload = _extract_agent_json_result(res_state, "kb_prompt_agent")
    return _normalize_tree_payload(payload, max_depth=max_depth)


@router.post("/generate")
async def generate_mindmap(
    request: Request,
    files: Optional[List[UploadFile]] = File(None),
    text: str = Form(""),
    chat_api_url: Optional[str] = Form(None),
    api_key: Optional[str] = Form(None),
    model: str = Form("gpt-5.4"),
    mindmap_style: str = Form("default"),
    max_depth: int = Form(3),
    language: str = Form("zh"),
    user: Optional[AuthUser] = Depends(get_optional_user),
):
    user = _require_or_resolve_user(user)
    file_items = files or []
    text_value = (text or "").strip()

    if not file_items and not text_value:
        raise HTTPException(status_code=400, detail="Please provide at least one file or some text content")

    resolved_api_url, resolved_api_key = resolve_llm_credentials(chat_api_url, api_key, scope="kb")
    owner = _owner_slug(user)
    run_id = f"{int(time.time())}_{uuid.uuid4().hex[:8]}"
    run_dir = (_mindmap_root_for_user(user) / run_id).resolve()
    inputs_dir = run_dir / "inputs"
    inputs_dir.mkdir(parents=True, exist_ok=True)

    local_paths: List[str] = []

    for uploaded in file_items:
        filename = Path(uploaded.filename or f"upload_{uuid.uuid4().hex}.txt").name
        suffix = Path(filename).suffix.lower()
        if suffix and suffix not in _ALLOWED_SUFFIXES:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {filename}")
        target = inputs_dir / filename
        target.write_bytes(await uploaded.read())
        local_paths.append(str(target))

    if text_value:
        prompt_path = inputs_dir / "notes.txt"
        prompt_path.write_text(text_value, encoding="utf-8")
        local_paths.append(str(prompt_path))

    parsed_files = await _parse_files(local_paths)
    tree_payload = await _generate_tree(
        text_blocks=parsed_files,
        chat_api_url=resolved_api_url,
        api_key=resolved_api_key,
        model=(model or "gpt-5.4").strip() or "gpt-5.4",
        max_depth=max(2, min(int(max_depth or 3), 6)),
        language=(language or "zh").strip() or "zh",
        style=(mindmap_style or "default").strip() or "default",
    )
    tree_root = tree_payload["root"]
    charge_info = estimate_mindmap_points(
        count_tree_nodes(tree_root),
        get_tree_depth(tree_root),
    )

    result_dir = run_dir.resolve()
    result_dir.mkdir(parents=True, exist_ok=True)
    mindmap_file = result_dir / "mindmap.json"
    payload_to_save = {
        "owner": owner,
        "style": (mindmap_style or "default").strip() or "default",
        "language": (language or "zh").strip() or "zh",
        "model": (model or "gpt-5.4").strip() or "gpt-5.4",
        "source_count": len(local_paths),
        "estimated_points": charge_info["points"],
        "billing": charge_info,
        "generated_at": int(time.time()),
        **tree_payload,
    }
    mindmap_file.write_text(json.dumps(payload_to_save, ensure_ascii=False, indent=2), encoding="utf-8")

    billing_service = BillingService()
    try:
        billing_service.consume_workflow(
            workflow_type="kb_mindmap",
            user=user,
            guest_id=getattr(request.state, "guest_id", None),
            amount=int(charge_info["points"]),
            event_key=f"mindmap_{run_id}",
        )
    except HTTPException:
        shutil.rmtree(result_dir, ignore_errors=True)
        raise

    return {
        "success": True,
        "result_path": _to_outputs_url(str(result_dir), request),
        "mindmap_path": _to_outputs_url(str(mindmap_file), request),
        "tree": tree_root,
        "highlights": tree_payload["highlights"],
        "input_count": len(local_paths),
        "estimated_points": charge_info["points"],
        "billing": charge_info,
        "pricing_summary": f"{charge_info['node_count']} nodes / depth {charge_info['depth']} = {charge_info['points']} points",
        "output_file_id": f"mindmap_{run_id}",
    }


@router.post("/save")
async def save_mindmap(
    request: Request,
    file_url: str = Body(..., embed=True),
    tree: Dict[str, Any] = Body(..., embed=True),
    highlights: Optional[List[str]] = Body(None, embed=True),
    user: Optional[AuthUser] = Depends(get_optional_user),
):
    _require_or_resolve_user(user)
    local_path = _resolve_owned_mindmap_path(file_url, user)
    if local_path.suffix.lower() not in {".json"}:
        raise HTTPException(status_code=400, detail="Invalid mindmap file type")

    normalized = _normalize_tree_payload({"root": tree, "highlights": highlights or []}, max_depth=6)
    local_path.parent.mkdir(parents=True, exist_ok=True)
    local_path.write_text(json.dumps(normalized, ensure_ascii=False, indent=2), encoding="utf-8")
    return {
        "success": True,
        "mindmap_path": _to_outputs_url(str(local_path), request),
        "tree": normalized["root"],
        "highlights": normalized["highlights"],
    }
