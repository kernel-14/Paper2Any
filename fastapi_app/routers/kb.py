import json as _json
import os
import re
import shutil
import subprocess
import time
import uuid
import zipfile
from pathlib import Path
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Body, Depends
from typing import Optional, List, Dict, Any

import fitz  # PyMuPDF

from dataflow_agent.state import IntelligentQARequest, IntelligentQAState, KBPodcastRequest, KBPodcastState, KBMindMapRequest, KBMindMapState
from dataflow_agent.utils import get_project_root
from dataflow_agent.workflow import run_workflow
from dataflow_agent.logger import get_logger
from fastapi_app.config import settings
from fastapi_app.dependencies import AuthUser, get_current_user
from fastapi_app.services.managed_api_service import is_free_billing_mode, resolve_llm_credentials
from fastapi_app.schemas import Paper2PPTRequest, DeepResearchRequest, DeepResearchResponse, KBReportRequest, KBReportResponse
from fastapi_app.utils import (
    _to_outputs_url,
    ensure_outputs_subpath,
    get_outputs_root,
    resolve_outputs_path,
)

router = APIRouter(prefix="/kb", tags=["Knowledge Base"])
log = get_logger(__name__)

# Base directory for storing KB files
# Use absolute path as requested by user or relative to project root
# We will use relative path 'outputs/kb_data' which resolves to that in the current workspace
KB_BASE_DIR = Path("outputs/kb_data")

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".mp4"}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}

NOTEBOOKS_DIR = KB_BASE_DIR / "_notebooks"

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


# ---------------------------------------------------------------------------
# Notebook helpers
# ---------------------------------------------------------------------------

def _safe_dirname(name: str) -> str:
    """Turn a notebook name into a filesystem-safe directory name.
    Keeps CJK characters, alphanumerics, hyphens and underscores.
    Collapses whitespace to '_' and strips dangerous chars."""
    if not name:
        return "unnamed"
    # Replace path separators and other dangerous chars
    cleaned = re.sub(r'[/\\:*?"<>|]', "", name)
    # Collapse whitespace to single underscore
    cleaned = re.sub(r"\s+", "_", cleaned).strip("_.")
    return cleaned or "unnamed"


def _canonical_user_id(user: AuthUser) -> str:
    user_id = (user.id or "").strip()
    if not user_id:
        raise HTTPException(status_code=400, detail="Authenticated user id is required")
    return user_id


def _canonical_user_email(user: AuthUser) -> str:
    email = (user.email or user.id or "").strip()
    if not email:
        raise HTTPException(status_code=400, detail="Authenticated user email is required")
    return email


def _resolve_kb_identity(user: AuthUser) -> tuple[str, str]:
    return _canonical_user_email(user), _canonical_user_id(user)


def _allowed_user_output_roots(user: AuthUser) -> List[Path]:
    email = _canonical_user_email(user)
    outputs_root = get_outputs_root()
    return [
        (outputs_root / "kb_data" / email).resolve(),
        (outputs_root / "kb_outputs" / email).resolve(),
        (outputs_root / "kb_exports" / email).resolve(),
    ]


def _ensure_path_within_roots(path: Path, roots: List[Path]) -> Path:
    resolved = path.resolve()
    for root in roots:
        try:
            resolved.relative_to(root)
            return resolved
        except ValueError:
            continue
    raise HTTPException(status_code=403, detail="Path does not belong to the authenticated user")


def _resolve_user_owned_output_path(path_or_url: str, user: AuthUser) -> Path:
    resolved = resolve_outputs_path(path_or_url, must_exist=False)
    return _ensure_path_within_roots(resolved, _allowed_user_output_roots(user))


def _notebooks_file(user_id: str) -> Path:
    """Return the path to the local notebooks JSON for a user."""
    NOTEBOOKS_DIR.mkdir(parents=True, exist_ok=True)
    return NOTEBOOKS_DIR / f"{user_id or 'default'}.json"


def _load_notebooks(user_id: str) -> List[Dict[str, Any]]:
    path = _notebooks_file(user_id)
    if not path.exists():
        return []
    try:
        return _json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return []


def _save_notebooks(user_id: str, notebooks: List[Dict[str, Any]]) -> None:
    path = _notebooks_file(user_id)
    path.write_text(_json.dumps(notebooks, ensure_ascii=False, indent=2), encoding="utf-8")


def _ensure_unique_dirname(user_id: str, base_name: str) -> str:
    """If base_name already used by another notebook, append _2, _3, etc."""
    notebooks = _load_notebooks(user_id)
    existing = {nb.get("dir_name") for nb in notebooks if nb.get("dir_name")}
    if base_name not in existing:
        return base_name
    seq = 2
    while f"{base_name}_{seq}" in existing:
        seq += 1
    return f"{base_name}_{seq}"


def _create_notebook_local(user_id: str, name: str, description: str = "") -> Dict[str, Any]:
    notebooks = _load_notebooks(user_id)
    nb_id = str(uuid.uuid4())
    dir_name = _safe_dirname(name)
    dir_name = _ensure_unique_dirname(user_id, dir_name)
    new_nb = {
        "id": nb_id,
        "name": name,
        "dir_name": dir_name,
        "description": description,
        "created_at": int(time.time()),
    }
    notebooks.append(new_nb)
    _save_notebooks(user_id, notebooks)
    return new_nb


def _get_notebook_local(user_id: str, notebook_id: str) -> Optional[Dict[str, Any]]:
    for nb in _load_notebooks(user_id):
        if nb.get("id") == notebook_id:
            return nb
    return None


def _get_notebook_dir_name(user_id: str, notebook_id: str) -> str:
    """Return the persisted dir_name for a notebook, or fall back to notebook_id."""
    nb = _get_notebook_local(user_id, notebook_id)
    if nb and nb.get("dir_name"):
        return nb["dir_name"]
    return notebook_id


def _notebook_base_dir(email: str, notebook_id: str, user_id: str = "default") -> Path:
    """kb_data/{email}/{notebook_dir_name}/
    Falls back to old path (using notebook_id) if new path doesn't exist."""
    dir_name = _get_notebook_dir_name(user_id, notebook_id)
    new_path = KB_BASE_DIR / email / dir_name
    if new_path.exists() or dir_name != notebook_id:
        return new_path
    # Fallback: old path using notebook_id directly
    old_path = KB_BASE_DIR / email / notebook_id
    if old_path.exists():
        return old_path
    # Neither exists yet — use new path
    return new_path


def _notebook_sources_dir(email: str, notebook_id: str, user_id: str = "default") -> Path:
    """kb_data/{email}/{notebook_name}/sources/"""
    return _notebook_base_dir(email, notebook_id, user_id) / "sources"


def _mineru_dir(email: str, notebook_id: str, user_id: str = "default") -> Path:
    """kb_data/{email}/{notebook_name}/mineru/"""
    return _notebook_base_dir(email, notebook_id, user_id) / "mineru"


def _vector_store_dir(email: str, notebook_id: str, user_id: str = "default") -> Path:
    """kb_data/{email}/{notebook_name}/vector_store/"""
    return _notebook_base_dir(email, notebook_id, user_id) / "vector_store"


def _generated_dir(
    email: str, notebook_id: str, output_type: str, user_id: str = "default"
) -> Path:
    """kb_data/{email}/{notebook_name}/generated/{output_type}/{next_seq}/
    Auto-increments the sequence number."""
    base = _notebook_base_dir(email, notebook_id, user_id) / "generated" / output_type
    base.mkdir(parents=True, exist_ok=True)
    existing = sorted(
        (int(d.name) for d in base.iterdir() if d.is_dir() and d.name.isdigit()),
        reverse=True,
    )
    next_seq = (existing[0] + 1) if existing else 1
    out = base / str(next_seq)
    out.mkdir(parents=True, exist_ok=True)
    return out


def _extract_text_result(state, role_name: str) -> str:
    try:
        result = getattr(state, "agent_results", {}).get(role_name, {}).get("results", {})
        if isinstance(result, dict):
            return result.get("text") or result.get("raw") or result.get("content") or ""
        if isinstance(result, str):
            return result
    except Exception:
        return ""
    return ""


def _extract_file_text(path: Path) -> str:
    if not path.exists():
        return ""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            doc = fitz.open(path)
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            return text
        if suffix in {".docx", ".doc"}:
            if Document is None:
                return ""
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        if suffix in {".pptx", ".ppt"}:
            if Presentation is None:
                return ""
            prs = Presentation(path)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        if suffix in {".md", ".txt"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        return ""
    except Exception:
        return ""


def _build_text_context(file_paths: List[str], max_chars: int = 60000) -> str:
    if not file_paths:
        return ""
    chunks: List[str] = []
    for raw in file_paths:
        try:
            local_path = _resolve_local_path(raw)
            _ensure_under_outputs(local_path)
            text = _extract_file_text(local_path)
            if not text:
                continue
            chunks.append(f"=== {local_path.name} ===\n{text}")
        except Exception:
            continue

    combined = "\n\n".join(chunks)
    if len(combined) > max_chars:
        combined = combined[:max_chars] + "\n\n[...content truncated]"
    return combined


def _get_deepresearch_service() -> "KBDeepResearchService":
    from fastapi_app.services.kb_deepresearch_service import KBDeepResearchService

    return KBDeepResearchService()


def _get_report_service() -> "KBReportService":
    from fastapi_app.services.kb_report_service import KBReportService

    return KBReportService()


def _safe_zip_stem(name: str) -> str:
    if not name:
        return "knowledge_base"
    cleaned = re.sub(r"[^A-Za-z0-9_-]+", "_", name.strip())
    return cleaned.strip("_") or "knowledge_base"


def _ensure_under_outputs(path: Path) -> None:
    ensure_outputs_subpath(path)


def _resolve_local_path(path_or_url: str) -> Path:
    return resolve_outputs_path(path_or_url, must_exist=False)


def _convert_to_pdf(input_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_path)
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise HTTPException(status_code=500, detail=f"PDF conversion failed for {input_path.name}")
    return pdf_path


def _merge_pdfs(pdf_paths: List[Path], output_path: Path) -> Path:
    if not pdf_paths:
        raise HTTPException(status_code=400, detail="No PDF files to merge")
    merged = fitz.open()
    for pdf in pdf_paths:
        with fitz.open(pdf) as src:
            merged.insert_pdf(src)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    merged.save(output_path)
    merged.close()
    return output_path


def _append_images_to_pptx(pptx_path: Path, image_paths: List[Path]) -> None:
    try:
        from pptx import Presentation
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"python-pptx not available: {e}")

    prs = Presentation(str(pptx_path))
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height
        )
    prs.save(str(pptx_path))

@router.post("/notebooks/create")
async def create_notebook(
    name: str = Body(..., embed=True),
    description: str = Body("", embed=True),
    user: AuthUser = Depends(get_current_user),
):
    """Create a new notebook and persist its dir_name mapping."""
    if not name or not name.strip():
        raise HTTPException(status_code=400, detail="Notebook name is required")
    nb = _create_notebook_local(_canonical_user_id(user), name.strip(), description)
    return {"success": True, "notebook": nb}


@router.get("/notebooks/list")
async def list_notebooks(user: AuthUser = Depends(get_current_user)):
    """List all notebooks for a user."""
    return {"success": True, "notebooks": _load_notebooks(_canonical_user_id(user))}


@router.get("/notebooks/{notebook_id}")
async def get_notebook(notebook_id: str, user: AuthUser = Depends(get_current_user)):
    """Get a single notebook by id."""
    nb = _get_notebook_local(_canonical_user_id(user), notebook_id)
    if not nb:
        raise HTTPException(status_code=404, detail="Notebook not found")
    return {"success": True, "notebook": nb}


@router.post("/notebooks/{notebook_id}/list-outputs")
async def list_notebook_outputs(
    notebook_id: str,
    email: Optional[str] = Body(None, embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    user: AuthUser = Depends(get_current_user),
):
    """Scan the generated/ directory of a notebook and return all outputs."""
    email, user_id = _resolve_kb_identity(user)
    base = _notebook_base_dir(email, notebook_id, user_id) / "generated"
    if not base.exists():
        return {"success": True, "outputs": {}}

    outputs: Dict[str, List[Dict[str, Any]]] = {}
    for type_dir in sorted(base.iterdir()):
        if not type_dir.is_dir():
            continue
        items: List[Dict[str, Any]] = []
        for seq_dir in sorted(type_dir.iterdir(), key=lambda d: d.name):
            if not seq_dir.is_dir() or not seq_dir.name.isdigit():
                continue
            files = []
            for f in seq_dir.rglob("*"):
                if f.is_file():
                    files.append({
                        "name": f.name,
                        "path": _to_outputs_url(str(f)),
                        "size": f.stat().st_size,
                    })
            items.append({"seq": int(seq_dir.name), "files": files})
        outputs[type_dir.name] = items
    return {"success": True, "outputs": outputs}


@router.post("/upload")
async def upload_kb_file(
    file: UploadFile = File(...),
    email: Optional[str] = Form(None),
    user_id: Optional[str] = Form(None),
    notebook_id: Optional[str] = Form(None),
    user: AuthUser = Depends(get_current_user),
):
    """
    Upload a file to the user's knowledge base directory.
    When notebook_id is provided: outputs/kb_data/{email}/{notebook_name}/sources/{filename}
    Legacy (no notebook_id): outputs/kb_data/{email}/{filename}
    """
    email, user_id = _resolve_kb_identity(user)

    # Validate file extension
    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type: {file_ext}. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"
        )

    try:
        # Determine target directory
        if notebook_id:
            user_dir = _notebook_sources_dir(email, notebook_id, user_id)
        else:
            user_dir = KB_BASE_DIR / email
        user_dir.mkdir(parents=True, exist_ok=True)

        # Secure filename (simple version)
        filename = file.filename
        if not filename:
            filename = f"unnamed_{user_id}"

        # Avoid path traversal
        filename = os.path.basename(filename)

        file_path = user_dir / filename

        # Save file
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        static_path = _to_outputs_url(str(file_path))

        return {
            "success": True,
            "filename": filename,
            "file_size": os.path.getsize(file_path),
            "storage_path": str(file_path),
            "static_url": static_path,
            "file_type": file.content_type
        }

    except Exception as e:
        log.error(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/delete")
async def delete_kb_file(
    storage_path: str = Form(...),
    user: AuthUser = Depends(get_current_user),
):
    """
    Delete a file from the physical storage.
    """
    try:
        target_path = _resolve_user_owned_output_path(storage_path, user)

        if target_path.exists() and target_path.is_file():
            os.remove(target_path)
            return {"success": True, "message": "File deleted"}
        else:
            return {"success": False, "message": "File not found"}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/delete-batch")
async def delete_kb_files_batch(
    storage_paths: List[str] = Body(..., embed=True),
    user: AuthUser = Depends(get_current_user),
):
    """
    Delete multiple files from physical storage.
    """
    if not storage_paths:
        raise HTTPException(status_code=400, detail="No storage paths provided")

    deleted = 0
    skipped: List[str] = []
    errors: List[str] = []

    for raw in storage_paths:
        try:
            local_path = _resolve_user_owned_output_path(raw, user)
            if local_path.exists() and local_path.is_file():
                local_path.unlink()
                deleted += 1
            else:
                skipped.append(raw)
        except Exception as e:
            errors.append(f"{raw}: {e}")

    return {
        "success": True,
        "deleted": deleted,
        "skipped": skipped,
        "errors": errors,
    }


@router.post("/export-zip")
async def export_kb_zip(
    files: List[str] = Body(..., embed=True),
    email: Optional[str] = Body(None, embed=True),
    kb_name: Optional[str] = Body(None, embed=True),
    include_root_dir: bool = Body(True, embed=True),
    user: AuthUser = Depends(get_current_user),
):
    """
    Export a list of KB files into a zip archive.
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    user_email = _canonical_user_email(user)
    outputs_root = get_outputs_root()
    export_root = outputs_root / "kb_exports" / user_email
    export_root.mkdir(parents=True, exist_ok=True)

    ts = int(time.time())
    zip_stem = _safe_zip_stem(kb_name or "knowledge_base")
    zip_path = export_root / f"{zip_stem}_{ts}.zip"

    used_names: Dict[str, int] = {}
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for raw in files:
            try:
                if not raw:
                    continue
                local_path = _resolve_user_owned_output_path(raw, user)
                if not local_path.exists() or not local_path.is_file():
                    continue
                base_name = local_path.name
                if base_name in used_names:
                    used_names[base_name] += 1
                    stem = local_path.stem
                    suffix = local_path.suffix
                    base_name = f"{stem}_{used_names[base_name]}{suffix}"
                else:
                    used_names[base_name] = 0

                arcname = base_name
                if include_root_dir:
                    arcname = f"{zip_stem}/{base_name}"
                zf.write(local_path, arcname)
            except Exception:
                continue

    if not zip_path.exists():
        raise HTTPException(status_code=500, detail="Failed to create zip")

    return {
        "success": True,
        "zip_path": _to_outputs_url(str(zip_path)),
        "count": len(files),
    }

@router.post("/chat")
async def chat_with_kb(
    files: List[str] = Body(..., embed=True),
    query: str = Body(..., embed=True),
    history: List[Dict[str, str]] = Body([], embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body(settings.KB_CHAT_MODEL, embed=True),
    user: AuthUser = Depends(get_current_user),
):
    """
    Intelligent QA Chat
    """
    try:
        resolved_api_url, resolved_api_key = resolve_llm_credentials(api_url, api_key, scope="kb")
        # Normalize file paths (web path -> local absolute path)
        local_files = []
        for f in files:
            local_path = _resolve_user_owned_output_path(f, user)
            if local_path.exists():
                local_files.append(str(local_path))

        if not local_files:
            raise HTTPException(status_code=400, detail="No valid files found")

        # Construct Request
        req = IntelligentQARequest(
            files=local_files,
            query=query,
            history=history,
            chat_api_url=resolved_api_url or os.getenv("DF_API_URL"),
            api_key=resolved_api_key or os.getenv("DF_API_KEY"),
            model=model
        )
        
        state = IntelligentQAState(request=req)
        
        # Run workflow via registry (统一使用 run_workflow)
        result_state = await run_workflow("intelligent_qa", state)
        
        # graph.ainvoke returns the final state dict or state object depending on implementation.
        # LangGraph usually returns dict. But our GenericGraphBuilder wrapper might return state.
        # GenericGraphBuilder compile returns a compiled graph.
        # Let's check typical usage. usually await graph.ainvoke(state) returns dict.
        
        answer = ""
        file_analyses = []
        
        if isinstance(result_state, dict):
            answer = result_state.get("answer", "")
            file_analyses = result_state.get("file_analyses", [])
        else:
            answer = getattr(result_state, "answer", "")
            file_analyses = getattr(result_state, "file_analyses", [])
            
        return {
            "success": True,
            "answer": answer,
            "file_analyses": file_analyses
        }

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

#
# Knowledge-base workflow endpoints were moved to:
#   fastapi_app/routers/kb_workflows.py
# This file now focuses on notebook / file-management / chat responsibilities.
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
