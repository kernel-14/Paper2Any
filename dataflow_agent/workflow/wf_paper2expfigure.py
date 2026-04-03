"""
paper2expfigure workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
从 PDF 论文中提取表格并生成统计图的完整工作流

支持三种输入类型：
- PDF: 从 PDF 论文中提取表格 (完整流程)
- FIGURE: 直接输入表格图片 (跳过 PDF 解析和 MinerU)
- TEXT: 输入表格文本，先生成表格图片再处理

工作流程：
1. PDF → 图片 (pdf_to_images_node)
2. 图片 → MinerU 识别 (mineru_extract_node)
3. 提取表格数据 (table_extractor_node)
4. 提取论文核心思想 (paper_idea_extractor_node) - TEXT/FIGURE 模式跳过
5. 智能推荐图表类型和生成代码 (code_executor_node)
   - 调用 chart_type_recommender Agent 推荐图表类型
   - 调用 chart_code_generator Agent 生成 matplotlib 代码
   - 执行代码生成图表
"""

from __future__ import annotations
import os
import uuid
import json
from pathlib import Path
from typing import Dict, Any, List
import asyncio
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass, field
from functools import reduce
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from dataflow_agent.state import Paper2FigureState
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.workflow.registry import register
from dataflow_agent.agentroles import create_simple_agent
from dataflow_agent.agentroles.paper2any_agents.chart_type_recommender import create_chart_type_recommender
from dataflow_agent.agentroles.paper2any_agents.chart_code_generator import create_chart_code_generator
from dataflow_agent.toolkits.tool_manager import get_tool_manager
from dataflow_agent.logger import get_logger
from dataflow_agent.utils import (
    pdf_to_pil_images,
    extract_tables_from_mineru_results,
    extract_text_from_mineru_results,
    execute_matplotlib_code,
)
from dataflow_agent.toolkits.multimodaltool.mineru_tool import run_aio_two_step_extract
from dataflow_agent.toolkits.multimodaltool.req_img import generate_or_edit_and_save_image_async
from dataflow_agent.utils.request_credentials import (
    get_request_image_api_key,
    get_request_image_api_url,
)


log = get_logger(__name__)


@register("paper2expfigure")
def create_paper2expfigure_graph() -> GenericGraphBuilder:
    """
    Paper2ExpFigure Workflow: 从 PDF 提取表格并生成统计图
    
    支持三种输入模式：
    - PDF: state.paper_file (完整流程)
    - FIGURE: state.fig_draft_path (表格图片，跳过 PDF 解析)
    - TEXT: state.paper_idea (表格文本，先生成图片)
    
    命令: dfa run --wf paper2expfigure
    """
    builder = GenericGraphBuilder(
        state_model=Paper2FigureState,
        entry_point="_start_"
    )

    # ======================================================================
    # PRE-TOOLS: 为 Agent 提供输入数据
    # ======================================================================
    
    @builder.pre_tool("paper_content", "paper_idea_extractor")
    def _get_paper_content(state: Paper2FigureState) -> str:
        """
        从 MinerU 结果或 PDF 中提取文本内容，供 paper_idea_extractor 使用
        """
        # 优先从 MinerU 结果中提取
        log.critical("正在从 MinerU 结果中提取文本内容")
        if hasattr(state, 'temp_data') and 'mineru_items' in state.temp_data:
            mineru_items = state.temp_data.get('mineru_items', [])
            if mineru_items:
                text = extract_text_from_mineru_results(mineru_items, max_chars=15000)
                if text:
                    return f"Paper content extracted from PDF:\n\n{text}"
        
        # 如果没有 MinerU 结果，直接从 PDF 读取（回退方案）
        import fitz
        pdf_path = state.paper_file
        if not pdf_path or not os.path.exists(pdf_path):
            log.warning("paper_file 为空或不存在，无法读取 PDF 内容")
            return ""
        
        try:
            doc = fitz.open(pdf_path)
            text_parts = []
            # 读取前 10 页
            for page_idx in range(min(10, len(doc))):
                page = doc.load_page(page_idx)
                text_parts.append(page.get_text("text") or "")
            doc.close()
            
            content = "\n".join(text_parts).strip()
            log.info(f"[pre_tool] 从 PDF 直接提取了 {len(content)} 字符")
            return f"Paper content from PDF:\n\n{content[:15000]}"
        except Exception as e:
            log.error(f"读取 PDF 失败: {e}")
            return ""
    
    # ==============================================================
    # 路由函数：根据输入类型决定流程
    # ==============================================================
    
    def _route_by_input_type(state: Paper2FigureState) -> str:
        """
        根据 input_type 决定下一个节点
        """
        input_type = getattr(state, 'input_type', None) or getattr(state.request, 'input_type', 'PDF')
        input_type = str(input_type).upper()
        
        log.info(f"[_route_by_input_type] input_type = {input_type}")
        
        if input_type == "FIGURE":
            return "figure_input_node"
        elif input_type == "TEXT":
            return "text_to_table_image_node"
        else:  # 默认 PDF
            return "pdf_to_images_node"
    
    # ==============================================================
    # NODES: 工作流节点
    # ==============================================================
    
    def _start_(state: Paper2FigureState) -> Paper2FigureState:
        """起始节点：初始化"""
        # 确保 temp_data 存在
        if not hasattr(state, 'temp_data') or state.temp_data is None:
            state.temp_data = {}
        
        # 确保 result_path 存在且为绝对路径
        if not state.result_path:
            output_dir = f"./outputs/paper2expfigure_{uuid.uuid4().hex[:8]}"
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_path.absolute())
        else:
            # 转换为绝对路径
            state.result_path = str(Path(state.result_path).resolve())
            Path(state.result_path).mkdir(parents=True, exist_ok=True)
        
        # 从 request 同步 input_type 到 state
        if hasattr(state.request, 'input_type'):
            state.input_type = state.request.input_type
        
        log.info(f"[_start_] result_path = {state.result_path}")
        log.info(f"[_start_] input_type = {getattr(state, 'input_type', 'PDF')}")
        
        return state
    
    async def figure_input_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        FIGURE 模式入口节点：处理直接输入的表格图片
        
        输入：state.fig_draft_path (单个图片路径或逗号分隔的多个路径)
        输出：构造与 MinerU 兼容的数据结构，直接进入 table_extractor_node
        """
        log.info("[figure_input_node] 开始处理表格图片输入...")
        
        fig_path = state.fig_draft_path or ""
        if not fig_path:
            log.error("[figure_input_node] fig_draft_path 为空")
            return state
        
        # 支持多个图片路径（逗号分隔）
        image_paths = [p.strip() for p in fig_path.split(",") if p.strip()]
        
        output_path = Path(state.result_path)
        table_images_dir = output_path / "table_images"
        table_images_dir.mkdir(parents=True, exist_ok=True)
        
        # 构造表格数据结构（跳过 MinerU，直接构造 extracted_tables）
        tables = []
        valid_image_paths = []
        
        for idx, img_path in enumerate(image_paths):
            img_path = Path(img_path)
            if not img_path.exists():
                log.warning(f"[figure_input_node] 图片不存在: {img_path}")
                continue
            
            # 复制图片到输出目录
            table_id = f"table_{idx}"
            dest_path = table_images_dir / f"{table_id}.png"
            
            try:
                img = Image.open(img_path)
                img.save(dest_path)
                valid_image_paths.append(str(dest_path))
                
                # 构造表格信息（没有 MinerU 解析，headers/rows 为空）
                tables.append({
                    "table_id": table_id,
                    "headers": [],
                    "rows": [],
                    "caption": f"Table from image: {img_path.name}",
                    "bbox": [0, 0, 1, 1],
                    "content": "",
                    "image_path": str(dest_path),
                    "page_index": 0,
                    "page_number": 1,
                })
                
                log.info(f"[figure_input_node] 处理图片 {idx + 1}: {img_path} -> {dest_path}")
                
            except Exception as e:
                log.error(f"[figure_input_node] 处理图片失败 ({img_path}): {e}")
        
        state.temp_data['image_paths'] = valid_image_paths
        state.extracted_tables = tables
        
        # FIGURE 模式不需要提取 paper_idea，设置默认值
        if not state.paper_idea:
            state.paper_idea = "Direct table image input - no paper context available"
        
        log.info(f"[figure_input_node] 完成，共处理 {len(tables)} 个表格图片")
        return state
    
    async def text_to_table_image_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        TEXT 模式入口节点：将表格文本转换为表格图片
        
        输入：state.paper_idea (表格文本，支持 CSV/Markdown/纯文本/LaTeX 等格式)
        输出：通过 LLM 生成 matplotlib 代码渲染表格图片，支持多级表头等复杂结构
        
        支持多表格：自动识别并分割文本中的多个表格，按 table_0, table_1... 命名
        """
        from dataflow_agent.agentroles.paper2any_agents.table_text_renderer import (
            render_table_from_text,
            split_tables_from_text,
        )
        
        log.info("[text_to_table_image_node] 开始处理表格文本输入...")
        
        table_text = state.paper_idea or ""
        if not table_text:
            log.error("[text_to_table_image_node] paper_idea (表格文本) 为空")
            return state
        
        output_path = Path(state.result_path).resolve()
        table_images_dir = output_path / "table_images"
        table_images_dir.mkdir(parents=True, exist_ok=True)
        
        # 先分割多个表格
        log.info("[text_to_table_image_node] 分析文本中的表格...")
        table_segments = await split_tables_from_text(
            text=table_text,
            state=state,
            model_name=state.request.model or "gpt-4o",
        )
        
        log.info(f"[text_to_table_image_node] 识别到 {len(table_segments)} 个表格")
        
        tables = []
        valid_image_paths = []
        
        # 循环处理每个表格
        for idx, segment in enumerate(table_segments):
            table_id = f"table_{idx}"
            img_path = (table_images_dir / f"{table_id}.png").resolve()
            segment_text = segment.get("text", "")
            caption = segment.get("caption", "")
            
            if not segment_text.strip():
                log.warning(f"[text_to_table_image_node] 表格 {table_id} 文本为空，跳过")
                continue
            
            log.info(f"[text_to_table_image_node] 处理表格 {idx + 1}/{len(table_segments)}: {table_id}")
            
            try:
                # 使用 table_text_renderer agent 渲染表格
                success, parsed_data = await render_table_from_text(
                    table_text=segment_text,
                    output_path=img_path,
                    state=state,
                    title=caption,
                    model_name=state.request.model or "gpt-4o",
                )
                
                if success:
                    valid_image_paths.append(str(img_path))
                    
                    tables.append({
                        "table_id": table_id,
                        "headers": parsed_data.get("headers", []),
                        "rows": parsed_data.get("rows", []),
                        "caption": caption,
                        "bbox": [0, 0, 1, 1],
                        "content": segment_text,
                        "image_path": str(img_path),
                        "page_index": 0,
                        "page_number": 1,
                        "has_multi_level_header": parsed_data.get("has_multi_level_header", False),
                    })
                    
                    log.info(f"[text_to_table_image_node] 生成表格图片: {img_path}")
                else:
                    log.warning(f"[text_to_table_image_node] 表格 {table_id} 渲染失败")
                
            except Exception as e:
                log.error(f"[text_to_table_image_node] 生成表格图片失败: {e}")
                import traceback
                traceback.print_exc()
        
        state.temp_data['image_paths'] = valid_image_paths
        state.extracted_tables = tables
        
        # TEXT 模式保留原始文本作为 paper_idea 的补充
        state.paper_idea = f"Table data from text input:\n{table_text}"
        
        log.info(f"[text_to_table_image_node] 完成，共生成 {len(tables)} 个表格图片")
        return state
    
    async def pdf_to_images_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 1: PDF → 图片
        将 PDF 的每一页转换为 PIL Image 对象，保存到临时目录
        """
        pdf_path = Path(state.paper_file)
        if not pdf_path.exists():
            log.error(f"PDF 文件不存在: {pdf_path}")
            return state
        
        log.info(f"[pdf_to_images] 开始转换 PDF: {pdf_path}")
        
        # 转换 PDF 为图片
        images = pdf_to_pil_images(pdf_path, dpi=150)
        
        # 创建临时目录保存图片
        output_path = Path(state.result_path)
        images_dir = output_path / "images"
        images_dir.mkdir(exist_ok=True)
        
        # 保存图片
        image_paths = []
        for idx, img in enumerate(images):
            img_path = images_dir / f"page_{idx+1}.png"
            img.save(img_path)
            image_paths.append(str(img_path))
            log.info(f"[pdf_to_images] 保存第 {idx+1} 页: {img_path}")
        
        # 存储到 state（使用绝对路径）
        state.temp_data['image_paths'] = image_paths
        
        log.info(f"[pdf_to_images] 完成，共转换 {len(images)} 页")
        return state
    
    async def mineru_extract_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 2: MinerU 识别
        使用 MinerU HTTP API 识别图片中的文本和表格
        """
        image_paths = state.temp_data.get('image_paths', [])
        if not image_paths:
            log.warning("[mineru_extract] 没有图片路径，跳过")
            return state
        
        output_path = Path(state.result_path)
        mineru_dir = output_path / "mineru_results"
        mineru_dir.mkdir(exist_ok=True)
        
        port = state.mineru_port
        all_items = []
        
        # 对每一页图片执行 MinerU 识别
        for idx, img_path in enumerate(image_paths, 1):
            log.info(f"[mineru_extract] 处理图片 {idx}/{len(image_paths)}: {img_path}")
            try:
                # 使用 run_aio_two_step_extract 进行识别
                items = await run_aio_two_step_extract(
                    image_path=str(img_path),
                    port=port,
                )
                
                # items 是一个列表，直接扩展到 all_items
                if isinstance(items, list):
                    # 为每个 item 添加页面信息
                    for item in items:
                        item['page_index'] = idx - 1  # 0-based index
                        item['page_number'] = idx      # 1-based number
                    
                    all_items.extend(items)
                    log.info(f"[mineru_extract] 从 page_{idx} 提取了 {len(items)} 个元素")
                    
                    # 保存每页的识别结果为 JSON 文件（便于调试）
                    result_file = mineru_dir / f"page_{idx}_result.json"
                    with open(result_file, 'w', encoding='utf-8') as f:
                        json.dump(items, f, ensure_ascii=False, indent=2)
                    log.debug(f"[mineru_extract] 保存结果到: {result_file}")
                else:
                    log.warning(f"[mineru_extract] 返回结果不是列表: {type(items)}")
                    
            except Exception as e:
                log.error(f"[mineru_extract] MinerU 识别失败 (page_{idx}): {e}")
        
        # 存储结果
        state.temp_data['mineru_items'] = all_items
        
        # 保存所有结果的汇总文件
        if all_items:
            summary_file = mineru_dir / "all_results.json"
            with open(summary_file, 'w', encoding='utf-8') as f:
                json.dump(all_items, f, ensure_ascii=False, indent=2)
            log.info(f"[mineru_extract] 保存汇总结果到: {summary_file}")
        
        log.info(f"[mineru_extract] 完成，共提取 {len(all_items)} 个元素")
        return state
    
    async def table_extractor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 3: 提取表格
        从 MinerU 识别结果中提取表格数据，并保存表格区域图片
        
        注意：FIGURE/TEXT 模式会跳过此节点（已在入口节点处理）
        """
        # 如果已经有 extracted_tables（FIGURE/TEXT 模式），跳过
        if state.extracted_tables:
            log.info(f"[table_extractor] 已有 {len(state.extracted_tables)} 个表格，跳过提取")
            return state
        
        mineru_items = state.temp_data.get('mineru_items', [])
        if not mineru_items:
            log.warning("[table_extractor] 没有 MinerU 结果，跳过")
            return state
        
        log.info("[table_extractor] 开始提取表格...")
        tables = extract_tables_from_mineru_results(mineru_items, min_rows=2, min_cols=2)
        
        log.info(f"[table_extractor] 提取了 {len(tables)} 个表格")
        
        # 打印表格摘要
        for table in tables:
            log.info(f"  - {table['table_id']}: {len(table['headers'])} 列 x {len(table['rows'])} 行")
        
        # 保存表格区域图片
        if tables:
            output_path = Path(state.result_path)
            table_images_dir = output_path / "table_images"
            table_images_dir.mkdir(exist_ok=True)
            
            # 获取原始图片路径
            image_paths = state.temp_data.get('image_paths', [])
            
            # 为每个表格裁剪并保存图片
            saved_count = 0
            for item in mineru_items:
                if item.get('type') != 'table':
                    continue
                
                bbox = item.get('bbox', [])
                if len(bbox) != 4:
                    continue
                
                # 找到对应的表格（通过 bbox 匹配）
                table_match = None
                for table in tables:
                    if table.get('bbox') == bbox:
                        table_match = table
                        break
                
                if not table_match:
                    continue
                
                # 从 item 中直接获取页面索引（在 mineru_extract_node 中添加的）
                page_idx = item.get('page_index')
                
                if page_idx is not None and page_idx < len(image_paths):
                    try:
                        # 读取原始图片
                        img_path = image_paths[page_idx]
                        img = Image.open(img_path)
                        
                        # bbox 是归一化坐标 [x0, y0, x1, y1]，范围 0-1
                        img_width, img_height = img.size
                        x0 = int(bbox[0] * img_width)
                        y0 = int(bbox[1] * img_height)
                        x1 = int(bbox[2] * img_width)
                        y1 = int(bbox[3] * img_height)
                        
                        # 裁剪表格区域
                        table_img = img.crop((x0, y0, x1, y1))
                        
                        # 保存图片
                        table_id = table_match['table_id']
                        page_num = item.get('page_number', page_idx + 1)
                        table_img_path = table_images_dir / f"{table_id}_page{page_num}.png"
                        table_img.save(table_img_path)
                        
                        # 将图片路径添加到 table 信息中
                        table_match['image_path'] = str(table_img_path)
                        table_match['page_index'] = page_idx
                        table_match['page_number'] = page_num
                        
                        saved_count += 1
                        log.info(f"[table_extractor] 保存表格图片: {table_img_path}")
                        
                    except Exception as e:
                        log.error(f"[table_extractor] 裁剪表格图片失败 ({table_match.get('table_id', 'unknown')}): {e}")
                        import traceback
                        traceback.print_exc()
            
            log.info(f"[table_extractor] 共保存了 {saved_count} 个表格图片到: {table_images_dir}")
        
        state.extracted_tables = tables
        
        return state
    
    async def paper_idea_extractor(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 4: 提取论文核心思想
        调用 paper_idea_extractor Agent 从论文中提取核心思想
        
        注意：FIGURE/TEXT 模式会跳过此节点
        """
        # 检查是否需要跳过（FIGURE/TEXT 模式已设置 paper_idea）
        input_type = getattr(state, 'input_type', None) or getattr(state.request, 'input_type', 'PDF')
        if input_type in ['FIGURE', 'TEXT']:
            log.info(f"[paper_idea_extractor] {input_type} 模式，跳过论文思想提取")
            return state
        
        log.info("[paper_idea_extractor] 开始提取论文核心思想...")
        
        agent = create_simple_agent(
            name="paper_idea_extractor",
            model_name=getattr(state.request, "chart_model", "gpt-4o"),
            temperature=0.1,
            max_tokens=4096,
            parser_type="json",
        )
        
        state = await agent.execute(state=state)
        
        paper_idea = state.paper_idea or ""
        log.info(f"[paper_idea_extractor] 提取的核心思想长度: {len(paper_idea)} 字符")
        log.info(f"[paper_idea_extractor] 核心思想预览: {paper_idea[:200]}...")
        
        return state

    async def chart_type_recommender_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 5: 智能推荐图表类型
        调用 chart_type_recommender Agent 智能推荐图表类型
        """
        log.info("[chart_type_recommender] 开始智能推荐图表类型...")
        
        tables = state.extracted_tables
        if not tables:
            log.warning("[chart_type_recommender] 没有表格数据，跳过")
            return state
        
        image_paths = [t["image_path"] for t in tables if "image_path" in t]
        
        # 这里直接用asyncio原生实现了并行，后面可以考虑改成更加符合langgraph的实现，使用Send API
        
        @dataclass
        class ChartTypeRecommenderState:
            request: Any
            pre_tool_results: Dict[str, Any]
            table: Dict
            chart_type_recommender: Any = None
            agent_results: Dict = field(default_factory=dict)
            chart_configs: Dict = field(default_factory=dict) # 这里存储最终结果：{table_id: chart_config}
        
        async def task(ctr_state: ChartTypeRecommenderState):
            table = ctr_state.table
            try:
                if "image_path" not in table:
                    log.error(f"[chart_type_recommender] 表格缺少图片路径: {table}")
                    return {table["table_id"]: None}
                input_image = table["image_path"]
            
                vlm_config = {
                    "mode": "understanding",
                    "input_image": input_image
                }

                agent = create_chart_type_recommender(
                    tool_manager=get_tool_manager(),
                    model_name=getattr(state.request, "chart_model", "gpt-4o"),
                    temperature=0.1,
                    max_tokens=2048,
                    vlm_config=vlm_config
                )
            
                ctr_state = await agent.execute(state=ctr_state)
                
                result = ctr_state.chart_configs
                
                return result
            except Exception as e:
                log.error(f"[chart_type_recommender] 处理表格出错: {e}")
                return {table["table_id"]: None}
        
        
        # 手动为每个并行节点注入 pre_tool_results
        
        states = [
            ChartTypeRecommenderState(
                request=state.request,
                table=table,
                pre_tool_results={
                    "table_info": {"table_id": table["table_id"]},
                    "paper_idea": state.paper_idea
                },
            )
            for table in tables
        ]
        
        log.debug(f"states: {states}")
        
        tasks = [task(s) for s in states]
        results = await asyncio.gather(*tasks)
        # 过滤掉失败的节点的返回值
        results = [
            result
            for result in results if result is not None
            for key, value in result.items() if value is not None
        ]
        results = reduce(lambda x, y: {**x, **y}, results, {})  # 合并结果为一个大字典
        state.chart_configs = results
        
        return state
    
    
    async def code_executor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 6: 执行代码生成图表
        调用 chart_code_generator Agent 智能生成图表
        """
        log.info("[code_executor] 开始执行代码生成图表...")

        tables = state.extracted_tables
        if not tables:
            log.warning("[code_executor] 没有表格数据，跳过")
            return state

        log.info(f"[code_executor] 共有 {len(tables)} 个表格待处理")
        
        image_paths = [t["image_path"] for t in tables if "image_path" in t]
        
        output_path = Path(state.result_path).resolve()
        charts_dir = output_path / "charts"
        charts_dir.mkdir(exist_ok=True)
        
        # 创建中间结果目录
        intermediate_dir = output_path / "chart_intermediate"
        intermediate_dir.mkdir(exist_ok=True)
        
        generated_charts = []
        
        # 获取论文核心思想
        paper_idea = state.paper_idea or "No paper idea extracted"
        
        @dataclass
        class ChartCodeGeneratorState:
            request: Any
            table: Dict
            result_path: str
            pre_tool_results: Dict[str, Any] = field(default_factory=dict)
            chart_code_generator: Any = None
            agent_results: Dict = field(default_factory=dict)
            generated_codes: Dict[str, Dict[str, Any]] = field(default_factory=dict)   # 生成的代码列表
        
        async def task(ccg_state: ChartCodeGeneratorState):
            table = ccg_state.table
            table_id = table['table_id']
            caption = table.get('caption', '')
            chart_config = ccg_state.pre_tool_results.get("chart_config", {})
            
            log.info(f"[code_executor] 处理表格: {table_id}")
            
            try:
                input_image = table["image_path"]
        
                vlm_config = {
                    "mode": "understanding",
                    "input_image": input_image
                }
                
                # 调用 Agent
                chart_code_agent = create_chart_code_generator(
                    tool_manager=get_tool_manager(),
                    model_name=getattr(state.request, "chart_model", "gpt-4o"),
                    temperature=0.0,
                    max_tokens=4096,
                    vlm_config=vlm_config,
                )
                
                ccg_state = await chart_code_agent.execute(state=ccg_state)
                
                # 获取生成的代码
                if ccg_state.generated_codes:
                    code_entry = ccg_state.generated_codes[table_id]  # 获取刚生成的代码
                    code = code_entry.get('code', '')
                    description = code_entry.get('description', '')
                    log.info(f"[code_executor] 生成代码长度: {len(code)} 字符")
                    log.info(f"[code_executor] 代码描述: {description}")
                else:
                    log.error(f"[code_executor] chart_code_generator 未返回代码")
                    raise Exception(f"chart_code_generator 未返回代码")
                
                # 4. 保存中间结果
                intermediate_file = intermediate_dir / f"{table_id}_intermediate.json"
                intermediate_data = {
                    "table_id": table_id,
                    "timestamp": str(Path(ccg_state.result_path).name),
                    
                    # 表格数据
                    "table_data": {
                        "caption": caption,
                    },
                    
                    # Agent 推荐结果
                    "chart_config": chart_config,
                    
                    # 生成的代码
                    "generated_code": code,
                    "code_description": description,
                }
                
                with open(intermediate_file, 'w', encoding='utf-8') as f:
                    json.dump(intermediate_data, f, ensure_ascii=False, indent=2)
                log.info(f"[code_executor] 保存中间结果: {intermediate_file}")
                
                # 保存代码文件（便于查看和调试）
                code_file = intermediate_dir / f"{table_id}_code.py"
                with open(code_file, 'w', encoding='utf-8') as f:
                    f.write(code)
                log.debug(f"[code_executor] 保存代码文件: {code_file}")
                
                # 5. 执行代码生成图表
                # 需要将 output_path, headers, rows 注入到代码中
                chart_path = (charts_dir / f"{table_id}.png").resolve()
                
                # 构建完整的可执行代码
                exec_code = f"""
# Auto-generated code execution wrapper
output_path = {repr(str(chart_path))}

# Generated chart code
{code}
"""
                
                result = execute_matplotlib_code(
                    code=exec_code,
                    output_path=chart_path,
                    timeout=30,
                )
                
                if result['success']:
                    generated_charts.append(str(chart_path))
                    log.info(f"[code_executor] 生成图表: {chart_path}")
                    
                    # 更新中间结果，添加执行状态
                    intermediate_data["execution_result"] = {
                        "success": True,
                        "chart_path": str(chart_path),
                        "error": None
                    }
                else:
                    log.error(f"[code_executor] 生成图表失败 ({table_id}): {result['error']}")
                    
                    # 更新中间结果，添加错误信息
                    intermediate_data["execution_result"] = {
                        "success": False,
                        "chart_path": None,
                        "error": result['error']
                    }
                
                # 重新保存中间结果（包含执行结果）
                with open(intermediate_file, 'w', encoding='utf-8') as f:
                    json.dump(intermediate_data, f, ensure_ascii=False, indent=2)
                
                code = ccg_state.generated_codes if ccg_state.generated_codes else {}
                chart_path = {table_id: chart_path}
                return (code, chart_path)
            
            except Exception as e:
                log.error(f"[code_executor] 处理表格 {table_id} 时出错: {e}")
                import traceback
                log.error(f"[code_executor] 错误堆栈:\n{traceback.format_exc()}")
                return ({table_id: None}, {table_id: None})
        
        # 过滤掉不适合生成图表的表格，定义匿名函数封装复杂提取逻辑，提高代码可读性
        get_chart_config = lambda x: state.chart_configs.get(x.get("table_id"), {})
        is_suitable = lambda x: get_chart_config(x).get("is_suitable_for_chart", True)
        
        # Debug: Print paper_idea length instead of full content to avoid log truncation
        log.info(f"[code_executor] paper_idea length: {len(state.paper_idea) if state.paper_idea else 0} characters")
        
        states = [
            ChartCodeGeneratorState(
                request=state.request,
                table=table,
                result_path=state.result_path,
                pre_tool_results={
                    "paper_idea": state.paper_idea,
                    "chart_config": get_chart_config(table),
                    "table_caption": table.get("caption", ""),
                }
            )
            for table in tables if is_suitable(table)
        ]
        
        tasks = [task(s) for s in states]
        generated_results = await asyncio.gather(*tasks)
        generated_code = [result[0] for result in generated_results]
        generated_charts = [result[1] for result in generated_results]

        # 过滤掉失败节点的值 - 修复：正确过滤 None 值
        code_results = {}
        for code_dict in generated_code:
            if code_dict:
                for table_id, code_entry in code_dict.items():
                    if code_entry is not None:
                        code_results[table_id] = code_entry

        chart_results = {}
        for chart_dict in generated_charts:
            if chart_dict:
                for table_id, chart_path in chart_dict.items():
                    if chart_path is not None:
                        chart_results[table_id] = chart_path
        
        state.generated_code = code_results
        state.generated_charts = chart_results
        
        log.info(f"[code_executor] 完成，共生成 {len(generated_charts)} 个图表")
        log.info(f"[code_executor] 中间结果保存在: {intermediate_dir}")
        
        return state

    async def post_stylize_node(state: Paper2FigureState) -> Paper2FigureState:
        """调用Nano Banana模型对生成的图表进行风格化，更美观"""
        log.info(f"[post_stylize] 开始")
        # 获取生成的图表路径列表，用于分发任务
        chart_paths = state.generated_charts
        
        if not chart_paths:
            log.warning("[post_stylize] 没有图表需要风格化，跳过")
            return state
        
        save_dir = Path(state.result_path) / Path("stylized_charts")
        save_dir.mkdir(parents=True, exist_ok=True)
        
        stylize_prompt = f"把这张统计图放大字体，并进行以 {state.request.style} 为主题的风格化，让它更美观，但需要保障数据的准确性，避免恶性的溢出和重叠"
        image_api_url = get_request_image_api_url(state.request)
        image_api_key = get_request_image_api_key(state.request)
        
        async def stylize_task(table_id: str, save_dir_path: str, chart_path: str):
            save_dir_p = Path(save_dir_path)
            chart_path_p = Path(chart_path)
            save_path = save_dir_p / chart_path_p.name.replace(".png", "_stylized.png")
            
            log.info(f"[post_stylize] 正在风格化图表: {table_id}")
            
            log.critical(f"image_path: {chart_path_p}")
            
            try:
                b64_result = await generate_or_edit_and_save_image_async(
                    prompt=stylize_prompt,
                    save_path=str(save_path),
                    api_url=image_api_url,
                    api_key=image_api_key,
                    model=state.request.gen_fig_model,
                    image_path=str(chart_path_p),
                    use_edit=True
                )
                
                log.info(f"[post_stylize] 图表风格化完成: {table_id}")
                return {table_id: [b64_result, save_path]}
            
            except Exception as e:
                log.error(f"[post_stylize] {table_id} 图表风格化出错: {e}")
                return {table_id: None}
        
        tasks = [stylize_task(table_id, str(save_dir), str(chart_path)) for table_id, chart_path in chart_paths.items()]
        results = await asyncio.gather(*tasks)
        # 过滤掉失败的图表
        results = [
            result 
            for result in results if result is not None 
            for table_id, stylize_result in result.items() if stylize_result is not None
        ]
        
        state.stylize_results = reduce(lambda x, y: {**x, **y}, results, {})
        
        log.info(f"[post_stylize] 完成")
        return state

    async def assemble_to_ppt(state: Paper2FigureState) -> Paper2FigureState:
        log.info(f"[assemble_to_ppt] 开始")
        
        try:
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            # 获取生成的图表路径
            chart_paths = state.generated_charts
            if not chart_paths:
                log.warning("[assemble_to_ppt] 没有生成的图表，跳过")
                return state
            
            # 获取风格化图片路径
            stylized_charts = state.stylize_results
            
            # 创建 PPT
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            def add_title_slide(title_text, subtitle_text=""):
                """添加标题页"""
                title_slide_layout = prs.slide_layouts[0]  # 标题布局
                slide = prs.slides.add_slide(title_slide_layout)
                
                # 设置标题
                title = slide.shapes.title
                title.text = title_text
                title.text_frame.paragraphs[0].font.size = Pt(44)
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 设置副标题
                if subtitle_text and len(slide.placeholders) > 1:
                    subtitle = slide.placeholders[1]
                    subtitle.text = subtitle_text
                    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
                    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                return slide
            
            def add_image_slide(image_path, title_text=""):
                """添加图片页"""
                blank_slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 添加标题（如果有）
                if title_text:
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = title_text
                    title_frame.paragraphs[0].font.size = Pt(24)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 获取图片尺寸
                img = Image.open(image_path)
                img_width, img_height = img.size
                
                # 计算适合幻灯片的尺寸（保持宽高比）
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # 为标题留出空间
                available_height = slide_height - Inches(1.2) if title_text else slide_height - Inches(0.5)
                max_width = slide_width - Inches(1)
                max_height = available_height
                
                # 计算缩放比例
                width_ratio = max_width / img_width
                height_ratio = max_height / img_height
                scale = min(width_ratio, height_ratio)
                
                # 计算最终尺寸
                final_width = int(img_width * scale)
                final_height = int(img_height * scale)
                
                # 居中放置
                left = (slide_width - final_width) / 2
                top_offset = Inches(1.2) if title_text else Inches(0.5)
                top = top_offset + (available_height - final_height) / 2
                
                # 添加图片
                slide.shapes.add_picture(
                    str(image_path),
                    left,
                    top,
                    width=final_width,
                    height=final_height
                )
                
                return slide
            
            # 1. 添加总标题页
            add_title_slide(
                "论文图表生成结果",
                "Paper2ExpFigure Workflow Results"
            )
            
            # 2. 添加原始图表部分标题页
            add_title_slide(
                "原始实验图表",
                "Original Experimental Charts"
            )
            
            # 3. 添加原始图表
            for table_id, chart_path in chart_paths.items():
                if not os.path.exists(chart_path):
                    log.warning(f"[assemble_to_ppt] 图表文件不存在: {chart_path}")
                    continue
                
                add_image_slide(chart_path, f"图表 {table_id}")
                log.info(f"[assemble_to_ppt] 添加原始图表到 PPT: {table_id}")
            
            # 4. 如果有风格化图片，添加风格化部分
            if stylized_charts:
                # 添加风格化图表部分标题页
                add_title_slide(
                    "风格化图表",
                    "Stylized Charts (Vintage Print Style)"
                )
                
                # 添加风格化图表
                for table_id, stylized_path in stylized_charts.items():
                    stylized_path = stylized_path[1]
                    if not os.path.exists(stylized_path):
                        log.warning(f"[assemble_to_ppt] 风格化图表文件不存在: {stylized_path}")
                        continue
                    
                    add_image_slide(stylized_path, f"风格化图表 {table_id}")
                    log.info(f"[assemble_to_ppt] 添加风格化图表到 PPT: {table_id}")
            
            # 保存 PPT
            output_path = Path(state.result_path)
            ppt_path = output_path / "generated_charts.pptx"
            prs.save(str(ppt_path))
            
            total_slides = len(prs.slides)
            log.info(f"[assemble_to_ppt] PPT 已保存: {ppt_path}")
            log.info(f"[assemble_to_ppt] 共创建 {total_slides} 张幻灯片")
            log.info(f"[assemble_to_ppt] 包含 {len(chart_paths)} 个原始图表和 {len(stylized_charts)} 个风格化图表")
            
            state.ppt_path = str(ppt_path)
            
        except ImportError:
            log.error("[assemble_to_ppt] 缺少 python-pptx 库，请安装: pip install python-pptx")
        except Exception as e:
            log.error(f"[assemble_to_ppt] 生成 PPT 失败: {e}")
            import traceback
            traceback.print_exc()
        
        log.info(f"[assemble_to_ppt] 完成")
        return state


    # ==============================================================
    # 注册 nodes / edges
    # ==============================================================
    
    nodes = {
        "_start_": _start_,
        "pdf_to_images_node": pdf_to_images_node,
        "mineru_extract_node": mineru_extract_node,
        "figure_input_node": figure_input_node,
        "text_to_table_image_node": text_to_table_image_node,
        "table_extractor_node": table_extractor_node,
        "paper_idea_extractor": paper_idea_extractor,
        "chart_type_recommender_node": chart_type_recommender_node,
        "code_executor_node": code_executor_node,
        "post_stylize_node": post_stylize_node,
        "assemble_to_ppt": assemble_to_ppt,
        "_end_": lambda state: state,
    }
    
    # 边定义
    edges = [
        # PDF 流程
        ("pdf_to_images_node", "mineru_extract_node"),
        ("mineru_extract_node", "table_extractor_node"),
        
        # FIGURE 流程 - 直接到 chart_type_recommender
        ("figure_input_node", "chart_type_recommender_node"),
        
        # TEXT 流程 - 直接到 chart_type_recommender
        ("text_to_table_image_node", "chart_type_recommender_node"),
        
        # PDF 流程继续
        ("table_extractor_node", "paper_idea_extractor"),
        ("paper_idea_extractor", "chart_type_recommender_node"),
        
        # 公共流程
        ("chart_type_recommender_node", "code_executor_node"),
        ("code_executor_node", "post_stylize_node"),
        ("post_stylize_node", "assemble_to_ppt"),
        
        # 最终节点
        ("assemble_to_ppt", "_end_"),
    ]
    
    # 添加条件路由
    builder.add_nodes(nodes).add_edges(edges).add_conditional_edge("_start_", _route_by_input_type)
    
    return builder
