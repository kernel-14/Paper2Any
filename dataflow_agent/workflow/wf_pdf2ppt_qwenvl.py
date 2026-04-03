"""
pdf2ppt_qwenvl workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
基于 slides PDF，融合 VLM (Qwen-VL-OCR) 替代传统 PaddleOCR：
1. 将 PDF 每页渲染为 PNG
2. 对每页图片用 VLM (ImageTextBBoxAgent) 做文字识别与定位 (替代 PaddleOCR)
3. 对每页图片用 SAM3 做分组提示词分割（与 paper2drawio_sam3 一致）
4. AI 背景编辑：
   - 基于 VLM 提取的 bbox 生成 mask (或利用 VLM 调试阶段的 no_text 图)
   - 调用 Inpainting API：填补 mask 之后的白色区域，结合背景颜色做 Inpainting
5. 智能合并与 PPT 生成：
   - 结合 SAM3 (图标), VLM (文字) 结果生成 PPT。
"""

from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional
import copy

import cv2
import numpy as np
import fitz  # PyMuPDF
from PIL import Image

from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.state import Paper2FigureState 
from dataflow_agent.utils import get_project_root
from dataflow_agent.agentroles import create_vlm_agent

# Tools
from dataflow_agent.toolkits.multimodaltool.sam3_tool import (
    Sam3PredictRun,
    decode_sam3_mask,
    dedup_sam3_results_across_groups,
    filter_sam3_items_contained_by_images,
    get_sam3_client,
    run_sam3_predict_runs,
)
from dataflow_agent.toolkits.multimodaltool.bg_tool import local_tool_for_bg_remove, free_bg_rm_model
from dataflow_agent.toolkits.multimodaltool.req_img import generate_or_edit_and_save_image_async
from dataflow_agent.utils.request_credentials import (
    get_request_image_api_key,
    get_request_image_api_url,
)
from dataflow_agent.toolkits.multimodaltool.ocr_config import get_ocr_api_credentials
from dataflow_agent.toolkits.multimodaltool.ocr_utils import extract_bbox_items
from dataflow_agent.toolkits.multimodaltool import ppt_tool
from dataflow_agent.workflow.sam3_segment_hint import (
    dedupe_prompts,
    generate_sam3_segment_hints,
)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

log = get_logger(__name__)



# --- Adaptive background fill helpers ---
def _norm_bbox_to_px(bbox_n: List[float], w: int, h: int, pad: int = 0) -> Optional[List[int]]:
    try:
        y1n, x1n, y2n, x2n = bbox_n
        x1 = int(x1n * w) - pad
        y1 = int(y1n * h) - pad
        x2 = int(x2n * w) + pad
        y2 = int(y2n * h) + pad
        x1 = max(0, min(w, x1))
        y1 = max(0, min(h, y1))
        x2 = max(0, min(w, x2))
        y2 = max(0, min(h, y2))
        if x2 <= x1 or y2 <= y1:
            return None
        return [x1, y1, x2, y2]
    except Exception:
        return None

def _expand_bbox_norm(bbox_n: List[float], w: int, h: int, pad_px: int) -> Optional[List[float]]:
    try:
        y1n, x1n, y2n, x2n = bbox_n
        x1 = max(0.0, (x1n * w - pad_px) / float(w))
        y1 = max(0.0, (y1n * h - pad_px) / float(h))
        x2 = min(1.0, (x2n * w + pad_px) / float(w))
        y2 = min(1.0, (y2n * h + pad_px) / float(h))
        if x2 <= x1 or y2 <= y1:
            return None
        return [y1, x1, y2, x2]
    except Exception:
        return None

def _bbox_overlap_ratio(inner: List[int], outer: List[int]) -> float:
    x1 = max(inner[0], outer[0])
    y1 = max(inner[1], outer[1])
    x2 = min(inner[2], outer[2])
    y2 = min(inner[3], outer[3])
    inter = max(0, x2 - x1) * max(0, y2 - y1)
    area = max(1, (inner[2] - inner[0]) * (inner[3] - inner[1]))
    return inter / area

def _adaptive_fill_text_regions(
    img_path: str,
    text_boxes_norm: List[List[float]],
    output_path: str,
    *,
    skip_boxes_norm: Optional[List[List[float]]] = None,
    pad: int = 2,
    ring: int = 6,
    std_low: float = 8.0,
    std_mid: float = 20.0,
    max_samples: int = 5000,
) -> bool:
    img = cv2.imread(img_path, cv2.IMREAD_COLOR)
    if img is None:
        return False
    h, w = img.shape[:2]
    out = img.copy()
    skip_boxes_px: List[List[int]] = []
    for sb in skip_boxes_norm or []:
        px = _norm_bbox_to_px(sb, w, h, pad=0)
        if px:
            skip_boxes_px.append(px)

    mask_high = np.zeros((h, w), dtype=np.uint8)

    for bbox_n in text_boxes_norm:
        px = _norm_bbox_to_px(bbox_n, w, h, pad=pad)
        if not px:
            continue
        x1, y1, x2, y2 = px

        # Skip text inside image/table zones to avoid damaging embedded figures
        if any(_bbox_overlap_ratio(px, sb) >= 0.5 for sb in skip_boxes_px):
            continue

        xo1 = max(0, x1 - ring)
        yo1 = max(0, y1 - ring)
        xo2 = min(w, x2 + ring)
        yo2 = min(h, y2 + ring)

        roi = img[yo1:yo2, xo1:xo2]
        if roi.size == 0:
            continue
        mask = np.ones((yo2 - yo1, xo2 - xo1), dtype=np.uint8)
        mask[y1 - yo1:y2 - yo1, x1 - xo1:x2 - xo1] = 0
        ys, xs = np.where(mask == 1)
        if len(xs) < 16:
            # Fallback: use whole ROI stats if ring too small
            ring_pixels = roi.reshape(-1, 3)
            xs = np.random.randint(0, roi.shape[1], size=min(len(ring_pixels), max_samples))
            ys = np.random.randint(0, roi.shape[0], size=min(len(ring_pixels), max_samples))
            ring_pixels = roi[ys, xs]
            xs_full = xs + xo1
            ys_full = ys + yo1
        else:
            if len(xs) > max_samples:
                idx = np.random.choice(len(xs), max_samples, replace=False)
                xs = xs[idx]
                ys = ys[idx]
            ring_pixels = roi[ys, xs]
            xs_full = xs + xo1
            ys_full = ys + yo1

        if ring_pixels.size == 0:
            continue

        ring_gray = (
            0.114 * ring_pixels[:, 0].astype(np.float32)
            + 0.587 * ring_pixels[:, 1].astype(np.float32)
            + 0.299 * ring_pixels[:, 2].astype(np.float32)
        )
        std = float(np.std(ring_gray))

        if std < std_low:
            mean_color = ring_pixels.mean(axis=0).astype(np.uint8)
            out[y1:y2, x1:x2] = mean_color
        elif std < std_mid:
            A = np.stack([xs_full, ys_full, np.ones_like(xs_full)], axis=1).astype(np.float32)
            xg, yg = np.meshgrid(np.arange(x1, x2), np.arange(y1, y2))
            for c in range(3):
                coeffs, _, _, _ = np.linalg.lstsq(A, ring_pixels[:, c].astype(np.float32), rcond=None)
                pred = coeffs[0] * xg + coeffs[1] * yg + coeffs[2]
                out[y1:y2, x1:x2, c] = np.clip(pred, 0, 255).astype(np.uint8)
        else:
            mask_high[y1:y2, x1:x2] = 255

    if np.any(mask_high):
        inpainted = cv2.inpaint(img, mask_high, 3, cv2.INPAINT_TELEA)
        out[mask_high == 255] = inpainted[mask_high == 255]

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    return bool(cv2.imwrite(output_path, out))

SHAPE_PROMPT = [
    "rectangle",
    "rounded rectangle",
    "diamond",
    "ellipse",
]

ARROW_PROMPT = [
    "arrow",
    "connector",
]

IMAGE_PROMPT = [
    "icon",
    "symbol",
    "pictogram",
    "logo",
]

IMAGE_PROMPT_RECALL = [
    "illustration",
    "character",
]

BACKGROUND_PROMPT = [
    "panel",
    "container",
    "filled region",
    "background",
]

SAM3_GROUPS = {
    "shape": SHAPE_PROMPT,
    "arrow": ARROW_PROMPT,
    "image": IMAGE_PROMPT,
    "background": BACKGROUND_PROMPT,
}

SAM3_GROUP_CONFIG = {
    "shape": {"score_threshold": 0.5, "min_area": 200, "priority": 3},
    "arrow": {"score_threshold": 0.45, "min_area": 50, "priority": 4},
    "image": {"score_threshold": 0.5, "min_area": 100, "priority": 2},
    "background": {"score_threshold": 0.25, "min_area": 500, "priority": 1},
}

SAM3_IMAGE_RECALL_SCORE_THRESHOLD = 0.38
SAM3_IMAGE_RECALL_MIN_AREA_BASE = 40
SAM3_IMAGE_RECALL_MIN_AREA_RATIO = 0.00003
SAM3_IMAGE_RECALL_TRIGGER_MAX_IMAGES = 2
SAM3_DEDUP_IOU = 0.7
SAM3_ARROW_DEDUP_IOU = 0.85
SAM3_SHAPE_IMAGE_IOU = 0.6
MAX_IMAGE_BBOX_AREA_RATIO = 0.88


def _bbox_area_px(b: List[int]) -> int:
    return max(0, b[2] - b[0]) * max(0, b[3] - b[1])


def _resolve_icon_bg_remove_mode() -> str:
    raw = str(os.getenv("PAPER2PPT_ICON_BG_REMOVE_MODE", "auto")).strip().lower()
    if raw in {"", "auto"}:
        return "auto"
    if raw in {"off", "false", "0", "disable", "disabled", "none"}:
        return "off"
    if raw == "cpu":
        return "cpu"
    if raw == "cuda" or raw.startswith("cuda:"):
        return raw
    log.warning(f"[pdf2ppt_qwenvl][RMBG] unknown mode={raw}, fallback to auto")
    return "auto"


def _sam3_predict_groups(
    client: Any,
    image_path: str,
    extra_image_prompts: Optional[List[str]] = None,
) -> List[Dict[str, Any]]:
    image_area: Optional[int] = None
    try:
        with Image.open(image_path) as img:
            image_area = int(img.width * img.height)
    except Exception:
        image_area = None

    merged_image_prompts = dedupe_prompts(IMAGE_PROMPT + (extra_image_prompts or []))
    merged_recall_prompts = dedupe_prompts(IMAGE_PROMPT_RECALL + (extra_image_prompts or []))

    try:
        base_runs: List[Sam3PredictRun] = []
        for group, prompts in SAM3_GROUPS.items():
            cfg = SAM3_GROUP_CONFIG.get(group, {})
            group_prompts = merged_image_prompts if group == "image" else prompts
            base_runs.append(
                Sam3PredictRun(
                    group=group,
                    prompts=group_prompts,
                    score_threshold=cfg.get("score_threshold"),
                    min_area=cfg.get("min_area"),
                )
            )

        all_results = run_sam3_predict_runs(
            client=client,
            image_path=image_path,
            runs=base_runs,
        )
        all_results = dedup_sam3_results_across_groups(
            all_results,
            group_config=SAM3_GROUP_CONFIG,
            dedup_iou=SAM3_DEDUP_IOU,
            arrow_dedup_iou=SAM3_ARROW_DEDUP_IOU,
            shape_image_iou=SAM3_SHAPE_IMAGE_IOU,
        )
        all_results = filter_sam3_items_contained_by_images(
            all_results,
            image_groups=["image"],
            contain_threshold=0.85,
        )

        safe_image_area = image_area if image_area and image_area > 0 else 1
        core_image_hits = 0
        for item in all_results:
            if item.get("group") != "image":
                continue
            bbox = item.get("bbox")
            if not bbox or len(bbox) != 4:
                continue
            bbox_ratio = float(_bbox_area_px(bbox)) / float(safe_image_area)
            if bbox_ratio <= MAX_IMAGE_BBOX_AREA_RATIO:
                core_image_hits += 1

        should_recall = core_image_hits <= SAM3_IMAGE_RECALL_TRIGGER_MAX_IMAGES
        if should_recall and merged_recall_prompts:
            recall_min_area = SAM3_IMAGE_RECALL_MIN_AREA_BASE
            if image_area and image_area > 0:
                recall_min_area = max(
                    SAM3_IMAGE_RECALL_MIN_AREA_BASE,
                    int(image_area * SAM3_IMAGE_RECALL_MIN_AREA_RATIO),
                )

            recall_results = run_sam3_predict_runs(
                client=client,
                image_path=image_path,
                runs=[
                    Sam3PredictRun(
                        group="image",
                        prompts=merged_recall_prompts,
                        score_threshold=SAM3_IMAGE_RECALL_SCORE_THRESHOLD,
                        min_area=recall_min_area,
                    )
                ],
            )
            if recall_results:
                all_results.extend(recall_results)
                all_results = dedup_sam3_results_across_groups(
                    all_results,
                    group_config=SAM3_GROUP_CONFIG,
                    dedup_iou=SAM3_DEDUP_IOU,
                    arrow_dedup_iou=SAM3_ARROW_DEDUP_IOU,
                    shape_image_iou=SAM3_SHAPE_IMAGE_IOU,
                )
                all_results = filter_sam3_items_contained_by_images(
                    all_results,
                    image_groups=["image"],
                    contain_threshold=0.85,
                )

        return all_results
    except Exception as e:
        log.warning(f"[pdf2ppt_qwenvl][SAM3] grouped predict failed: {e}")
        return []


def _save_sam3_crop(
    image_bgr: np.ndarray,
    bbox: List[int],
    mask: Optional[np.ndarray],
    out_path: Path,
) -> bool:
    x1, y1, x2, y2 = bbox
    if x2 <= x1 or y2 <= y1:
        return False

    crop = image_bgr[y1:y2, x1:x2]
    if crop.size == 0:
        return False

    if mask is not None and mask.shape[:2] == image_bgr.shape[:2]:
        mask_u8 = (mask > 0).astype(np.uint8)
        alpha = (mask_u8[y1:y2, x1:x2] * 255).astype(np.uint8)
        rgba = np.dstack((crop, alpha))
        return bool(cv2.imwrite(str(out_path), rgba))

    return bool(cv2.imwrite(str(out_path), crop))

def get_closest_aspect_ratio(w: int, h: int) -> str:
    """
    计算最接近的合法 Gemini 比例
    """
    valid_ratios = ['1:1', '2:3', '3:2', '3:4', '4:3', '4:5', '5:4', '9:16', '16:9', '21:9']
    target_ratio = w / h
    
    best_ratio = '16:9' # default
    min_diff = float('inf')
    
    for r_str in valid_ratios:
        rw, rh = map(int, r_str.split(':'))
        curr_ratio = rw / rh
        diff = abs(target_ratio - curr_ratio)
        if diff < min_diff:
            min_diff = diff
            best_ratio = r_str
            
    return best_ratio

def _ensure_result_path(state: Paper2FigureState) -> str:
    raw = getattr(state, "result_path", None)
    if raw:
        return raw
    root = get_project_root()
    ts = int(__import__("time").time())
    base_dir = (root / "outputs" / "pdf2ppt_qwenvl" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path

def _process_single_sam_page(
    page_idx: int,
    img_path: str,
    base_dir: str,
    extra_image_prompts: Optional[List[str]] = None,
) -> Dict[str, Any]:
    log.info(f"[pdf2ppt_qwenvl][SAM3] processing page#{page_idx+1}: {img_path}")
    img_path_obj = Path(img_path)
    if not img_path_obj.exists():
        log.warning(f"[pdf2ppt_qwenvl][SAM3] page#{page_idx+1} image not found")
        return {"page_idx": page_idx, "layout_items": []}

    out_dir = Path(base_dir) / "layout_items" / f"page_{page_idx+1:03d}"
    out_dir.mkdir(parents=True, exist_ok=True)

    image_bgr = cv2.imread(str(img_path_obj), cv2.IMREAD_COLOR)
    if image_bgr is None:
        log.warning(f"[pdf2ppt_qwenvl][SAM3] page#{page_idx+1} read image failed")
        return {"page_idx": page_idx, "layout_items": []}

    h, w = image_bgr.shape[:2]
    client = get_sam3_client()
    if client is None:
        log.error("[pdf2ppt_qwenvl][SAM3] endpoints not configured")
        return {"page_idx": page_idx, "layout_items": []}

    layout_items: List[Dict[str, Any]] = []
    try:
        sam3_results = _sam3_predict_groups(
            client,
            str(img_path_obj),
            extra_image_prompts=extra_image_prompts,
        )
        for idx, item in enumerate(sam3_results):
            group = item.get("group") or ""
            if group == "background":
                continue

            bbox_raw = item.get("bbox")
            if not bbox_raw or len(bbox_raw) != 4:
                continue

            x1, y1, x2, y2 = [int(v) for v in bbox_raw]
            x1 = max(0, min(w, x1))
            x2 = max(0, min(w, x2))
            y1 = max(0, min(h, y1))
            y2 = max(0, min(h, y2))
            if x2 <= x1 or y2 <= y1:
                continue

            mask = decode_sam3_mask(item.get("mask") or {})
            crop_path = out_dir / f"sam3_{group}_{idx}.png"
            if not _save_sam3_crop(image_bgr, [x1, y1, x2, y2], mask, crop_path):
                continue

            layout_items.append(
                {
                    "bbox": [x1 / max(1, w), y1 / max(1, h), x2 / max(1, w), y2 / max(1, h)],
                    "bbox_px": [x1, y1, x2, y2],
                    "score": item.get("score"),
                    "area": max(0, (x2 - x1) * (y2 - y1)),
                    "prompt": item.get("prompt"),
                    "group": group,
                    "png_path": str(crop_path),
                    "type": "layout_box",
                }
            )
        log.info(
            f"[pdf2ppt_qwenvl][SAM3] page#{page_idx+1} returned {len(layout_items)} items "
            f"extra_image_prompts={extra_image_prompts or []}"
        )
    except Exception as e:
        log.error(f"[pdf2ppt_qwenvl][SAM3] page#{page_idx+1} failed: {e}")
        layout_items = []

    return {"page_idx": page_idx, "layout_items": layout_items}

def _run_sam_on_pages(
    image_paths: List[str],
    base_dir: str,
    extra_image_prompts_by_page: Optional[Dict[int, List[str]]] = None,
) -> List[Dict[str, Any]]:
    results: List[Dict[str, Any]] = []
    log.info(f"[pdf2ppt_qwenvl][SAM3] start, images={len(image_paths)}")
    for idx, path in enumerate(image_paths):
        try:
            results.append(
                _process_single_sam_page(
                    idx,
                    path,
                    base_dir,
                    extra_image_prompts=(extra_image_prompts_by_page or {}).get(idx) or [],
                )
            )
        except Exception as e:
            log.error(f"[pdf2ppt_qwenvl][SAM3] page#{idx+1} task exception: {e}")
            results.append({"page_idx": idx, "layout_items": []})

    return sorted(results, key=lambda x: x["page_idx"])

@register("pdf2ppt_qwenvl")
def create_pdf2ppt_qwenvl_graph() -> GenericGraphBuilder:
    """
    Workflow factory: dfa run --wf pdf2ppt_qwenvl
    """
    builder = GenericGraphBuilder(state_model=Paper2FigureState, entry_point="_start_")

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        _ensure_result_path(state)
        return state

    async def pdf_to_images_node(state: Paper2FigureState) -> Paper2FigureState:
        if state.request.input_type == "FIGURE":
            img_path = state.request.input_content
            log.info(f"[pdf2ppt_qwenvl] FIGURE mode: using input image {img_path}")
            if img_path and os.path.exists(img_path):
                state.slide_images = [img_path]
            else:
                log.error(f"[pdf2ppt_qwenvl] FIGURE mode: image not found {img_path}")
            return state

        pdf_path = getattr(state, "pdf_file", None)
        if not pdf_path:
            log.error("[pdf2ppt_qwenvl] state.pdf_file is empty")
            return state

        base_dir = Path(_ensure_result_path(state))
        img_dir = base_dir / "slides_png"
        image_paths = ppt_tool.pdf_to_images(pdf_path, str(img_dir))
        state.slide_images = image_paths
        return state

    # --- 新增：VLM 节点 (替代原 OCR) ---
    async def vlm_recognition_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        使用 VLM (ImageTextBBoxAgent) 提取文本和 bbox。
        结果写入 state.vlm_pages。
        """
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            log.warning("[pdf2ppt_qwenvl][vlm] no slide_images")
            state.vlm_pages = []
            return state

        async def _process_single_image(page_idx: int, img_path: str) -> Dict[str, Any]:
            try:
                # 显式传递 result_path，确保 agent 内部能访问
                temp_state = copy.copy(state)
                temp_state.result_path = state.result_path
                
                # Retry loop for VLM execution
                max_retries = 3
                bbox_res = []
                
                for attempt in range(max_retries):
                    try:
                        ocr_api_url, ocr_api_key = get_ocr_api_credentials()
                        if getattr(temp_state, "request", None):
                            temp_state.request = copy.copy(state.request)
                            temp_state.request.chat_api_url = ocr_api_url
                            temp_state.request.api_key = ocr_api_key
                            temp_state.request.chat_api_key = ocr_api_key

                        agent = create_vlm_agent(
                            name="ImageTextBBoxAgent",
                            model_name=getattr(state.request, "vlm_model", "qwen-vl-ocr-2025-11-20"),
                            chat_api_url=ocr_api_url,
                            max_tokens=4096,
                            vlm_mode="ocr",
                            additional_params={
                                "input_image": img_path
                            }
                        )

                        log.info(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} attempt {attempt+1}/{max_retries}...")
                        new_state = await agent.execute(temp_state)
                        raw_bbox_res = getattr(new_state, "bbox_result", [])
                        bbox_res = extract_bbox_items(raw_bbox_res)

                        if isinstance(raw_bbox_res, list):
                            break
                        if bbox_res:
                            break
                        log.warning(
                            f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} attempt {attempt+1} "
                            f"got invalid result: {type(raw_bbox_res)}"
                        )
                    except Exception as e:
                        log.warning(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} attempt {attempt+1} failed: {e}")
                        if attempt == max_retries - 1:
                            raise e
                        await asyncio.sleep(1)

                if not isinstance(bbox_res, list):
                    bbox_res = []
                
                # 修正 bbox 归一化 (0-1000 -> 0-1)
                # 并生成 "no_text" mask 图，供后续 Inpainting 使用
                processed_items = []
                
                # 读取原图尺寸
                pil_img = Image.open(img_path)
                w, h = pil_img.size
                
                # 准备生成 no_text 图
                mask_img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                
                VLM_SCALE = 1000.0
                
                for it in bbox_res:
                    # 处理 rotate_rect
                    if "rotate_rect" in it and "bbox" not in it:
                        try:
                            rr = it["rotate_rect"]
                            if isinstance(rr, list) and len(rr) == 5:
                                cx, cy, rw, rh, angle = rr
                                rect = ((float(cx), float(cy)), (float(rw), float(rh)), float(angle))
                                box = cv2.boxPoints(rect)
                                x_min = np.min(box[:, 0])
                                x_max = np.max(box[:, 0])
                                y_min = np.min(box[:, 1])
                                y_max = np.max(box[:, 1])
                                
                                it["bbox"] = [
                                    max(0.0, min(1.0, y_min / VLM_SCALE)),
                                    max(0.0, min(1.0, x_min / VLM_SCALE)),
                                    max(0.0, min(1.0, y_max / VLM_SCALE)),
                                    max(0.0, min(1.0, x_max / VLM_SCALE))
                                ]
                        except Exception:
                            pass
                    
                    if "bbox" in it:
                        processed_items.append(it)
                        # 在 mask_img 上将文字区域涂白
                        y1_n, x1_n, y2_n, x2_n = it["bbox"]
                        x1 = int(x1_n * w)
                        y1 = int(y1_n * h)
                        x2 = int(x2_n * w)
                        y2 = int(y2_n * h)
                        # 按文字高度扩张 mask，减少残影
                        box_h = max(1, y2 - y1)
                        pad = max(4, int(0.15 * box_h))
                        x1 = max(0, x1 - pad)
                        y1 = max(0, y1 - pad)
                        x2 = min(w, x2 + pad)
                        y2 = min(h, y2 + pad)
                        
                        cv2.rectangle(mask_img, (x1, y1), (x2, y2), (255, 255, 255), -1)

                # 保存 no_text 图片
                base_dir = Path(_ensure_result_path(state))
                debug_dir = base_dir / "vlm_debug"
                debug_dir.mkdir(parents=True, exist_ok=True)
                no_text_path = debug_dir / f"page_{page_idx+1:03d}_no_text.png"
                cv2.imwrite(str(no_text_path), mask_img)
                
                log.info(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} items={len(processed_items)}, saved mask to {no_text_path}")
                
                return {
                    "page_idx": page_idx,
                    "path": img_path,
                    "vlm_data": processed_items,
                    "no_text_path": str(no_text_path)
                }

            except Exception as e:
                log.error(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} failed: {e}")
                return {
                    "page_idx": page_idx,
                    "path": img_path,
                    "vlm_data": [],
                    "error": str(e)
                }

        tasks = [_process_single_image(i, p) for i, p in enumerate(image_paths)]
        results = await asyncio.gather(*tasks)
        state.vlm_pages = results
        return state

    async def slides_segment_hint_node(state: Paper2FigureState) -> Paper2FigureState:
        vlm_pages: List[Dict[str, Any]] = getattr(state, "vlm_pages", []) or []
        state.temp_data["pdf2ppt_sam3_segment_hints_by_page"] = {}
        state.temp_data["pdf2ppt_sam3_segment_hints_raw_by_page"] = {}
        if not vlm_pages:
            return state

        async def _process_single_page(pinfo: Dict[str, Any]) -> tuple[int, List[str], Any]:
            page_idx = int(pinfo.get("page_idx", 0))
            img_path = pinfo.get("path") or pinfo.get("sam_input_path") or ""
            if not img_path or not os.path.exists(img_path):
                return page_idx, [], {}
            try:
                hints, raw_result, _, _ = await generate_sam3_segment_hints(
                    state=state,
                    image_path=img_path,
                    text_blocks=pinfo.get("vlm_data", []) or [],
                    env_prefix="PAPER2PPT_SEGMENT_HINT",
                    base_image_prompts=IMAGE_PROMPT,
                    base_recall_prompts=IMAGE_PROMPT_RECALL,
                    extra_blocked_prompts=SHAPE_PROMPT + ARROW_PROMPT + BACKGROUND_PROMPT,
                    log_prefix=f"[pdf2ppt_qwenvl][segment_hint][page#{page_idx+1}]",
                )
                return page_idx, hints, raw_result
            except Exception as e:
                log.warning(f"[pdf2ppt_qwenvl][segment_hint] page#{page_idx+1} failed: {e}")
                return page_idx, [], {"error": str(e)}

        results = await asyncio.gather(*[_process_single_page(p) for p in vlm_pages])
        hints_by_page: Dict[int, List[str]] = {}
        raw_by_page: Dict[int, Any] = {}
        for page_idx, hints, raw_result in results:
            hints_by_page[page_idx] = hints
            raw_by_page[page_idx] = raw_result

        state.temp_data["pdf2ppt_sam3_segment_hints_by_page"] = hints_by_page
        state.temp_data["pdf2ppt_sam3_segment_hints_raw_by_page"] = raw_by_page
        for pinfo in vlm_pages:
            page_idx = int(pinfo.get("page_idx", 0))
            pinfo["sam3_segment_hints"] = hints_by_page.get(page_idx, [])
        return state

    async def slides_sam_node(state: Paper2FigureState) -> Paper2FigureState:
        """SAM3 图标分割"""
        vlm_pages: List[Dict[str, Any]] = getattr(state, "vlm_pages", []) or []
        slide_images: List[str] = getattr(state, "slide_images", []) or []
        image_paths: List[str] = []
        if vlm_pages and slide_images:
            page_dict = {p.get("page_idx", 0): p for p in vlm_pages}
            for idx, fallback in enumerate(slide_images):
                pinfo = page_dict.get(idx, {})
                candidates = [
                    pinfo.get("path"),
                    fallback,
                ]
                chosen = None
                for c in candidates:
                    if c and os.path.exists(c):
                        chosen = c
                        break
                if not chosen:
                    chosen = fallback
                if pinfo:
                    pinfo["sam_input_path"] = chosen
                image_paths.append(chosen)
        else:
            image_paths = slide_images
        log.info(f"[pdf2ppt_qwenvl][SAM3] start, images={len(image_paths)}")
        if not image_paths:
            return state
        base_dir = _ensure_result_path(state)
        extra_image_prompts_by_page: Dict[int, List[str]] = (
            state.temp_data.get("pdf2ppt_sam3_segment_hints_by_page", {}) or {}
        )
        try:
            sam_pages = await asyncio.to_thread(
                _run_sam_on_pages,
                image_paths,
                base_dir,
                extra_image_prompts_by_page,
            )
            log.info(f"[pdf2ppt_qwenvl][SAM3] done, pages={len(sam_pages)}")
        except Exception as e:
            log.error(f"[pdf2ppt_qwenvl][SAM3] processing failed: {e}")
            sam_pages = []
        state.sam_pages = sam_pages
        return state

    async def slides_layout_bg_remove_node(state: Paper2FigureState, sam_pages: List[Dict[str, Any]] = None) -> Paper2FigureState:
        """SAM 结果背景移除"""
        if sam_pages is None:
            sam_pages = getattr(state, "sam_pages", []) or []
        if not sam_pages:
            return state

        bg_remove_mode = _resolve_icon_bg_remove_mode()
        total_items = sum(len(p.get("layout_items", []) or []) for p in sam_pages)
        if bg_remove_mode == "off":
            for p in sam_pages:
                for it in p.get("layout_items", []) or []:
                    png_path = it.get("png_path")
                    if png_path and os.path.exists(png_path):
                        it["fg_png_path"] = png_path
            log.info(
                f"[pdf2ppt_qwenvl][RMBG] skipped icon bg remove: mode=off items={total_items}"
            )
            state.sam_pages = sam_pages
            return state

        base_dir = Path(_ensure_result_path(state))
        icons_dir = base_dir / "sam_icons"
        icons_dir.mkdir(parents=True, exist_ok=True)
        model_path = getattr(getattr(state, "request", None), "bg_rm_model", None)

        def _sync_bg_remove():
            processed = 0
            for p in sam_pages:
                page_idx = p.get("page_idx", 0)
                for it in p.get("layout_items", []):
                    png_path = it.get("png_path")
                    if not png_path or not os.path.exists(png_path): continue
                    
                    try:
                        original_stem = Path(png_path).stem
                        output_filename = f"page_{page_idx+1:03d}_{original_stem}_bg_removed.png"
                        output_path = icons_dir / output_filename
                        
                        req = {"image_path": png_path, "output_dir": str(icons_dir)}
                        if model_path: req["model_path"] = model_path
                        if bg_remove_mode != "auto": req["device"] = bg_remove_mode
                        
                        fg_path = local_tool_for_bg_remove(req)
                        
                        if fg_path and os.path.exists(fg_path):
                            fg_path_obj = Path(fg_path)
                            if fg_path_obj.name != output_filename:
                                new_fg_path = fg_path_obj.parent / output_filename
                                fg_path_obj.rename(new_fg_path)
                                fg_path = str(new_fg_path)
                            it["fg_png_path"] = fg_path
                        else:
                            it["fg_png_path"] = png_path
                        processed += 1
                    except Exception:
                        it["fg_png_path"] = png_path
            
            try:
                if model_path: free_bg_rm_model(model_path=model_path)
            except Exception: pass
            return processed

        log.info(
            f"[pdf2ppt_qwenvl][RMBG] start icon bg remove: mode={bg_remove_mode} items={total_items}"
        )
        await asyncio.to_thread(_sync_bg_remove)
        state.sam_pages = sam_pages
        return state

    async def slides_inpainting_node(state: Paper2FigureState) -> Paper2FigureState:
        """AI Inpainting: 填补文字 mask 区域"""
        vlm_pages = getattr(state, "vlm_pages", []) or []
        if not vlm_pages:
            return state

        base_dir = Path(_ensure_result_path(state))
        sam_pages = getattr(state, "sam_pages", []) or []
        sam_dict = {p.get("page_idx", 0): (p.get("layout_items", []) or []) for p in sam_pages}
        
        # API 配置
        req_cfg = getattr(state, "request", None) or {}
        if not isinstance(req_cfg, dict): req_cfg = req_cfg.__dict__ if hasattr(req_cfg, "__dict__") else {}
        api_key = get_request_image_api_key(getattr(state, "request", None))
        api_url = get_request_image_api_url(getattr(state, "request", None)) or "https://api.apiyi.com"
        model_name = req_cfg.get("gen_fig_model") or "gemini-3-pro-image-preview"
        is_comfly = "comfly" in str(api_url).lower()
        
        # 限制并发
        sem = asyncio.Semaphore(3)

        async def _call_image_api_with_retry(coro_factory, retries=3):
            for i in range(retries):
                try:
                    await coro_factory()
                    return True
                except Exception as e:
                    if i == retries - 1: log.error(f"Image API failed: {e}")
                    await asyncio.sleep(1)
            return False

        def _cleanup_bg_overlay(
            page_idx: int,
            pinfo: Dict[str, Any],
            bg_path: str,
            img_w: int,
            img_h: int,
        ) -> Optional[str]:
            if not bg_path or not os.path.exists(bg_path) or img_w <= 0 or img_h <= 0:
                return None

            bg_img = cv2.imread(bg_path, cv2.IMREAD_COLOR)
            if bg_img is None:
                return None

            if bg_img.shape[1] != img_w or bg_img.shape[0] != img_h:
                bg_img = cv2.resize(bg_img, (img_w, img_h), interpolation=cv2.INTER_CUBIC)

            overlap_mask = np.zeros((img_h, img_w), dtype=np.uint8)

            for item in pinfo.get("vlm_data", []) or []:
                bbox_n = item.get("bbox")
                if not bbox_n or len(bbox_n) != 4:
                    continue
                y1n, x1n, y2n, x2n = bbox_n
                x1 = int(x1n * img_w)
                y1 = int(y1n * img_h)
                x2 = int(x2n * img_w)
                y2 = int(y2n * img_h)
                box_h = max(1, y2 - y1)
                pad = max(4, int(0.15 * box_h))
                x1 = max(0, x1 - pad)
                y1 = max(0, y1 - pad)
                x2 = min(img_w, x2 + pad)
                y2 = min(img_h, y2 + pad)
                if x2 > x1 and y2 > y1:
                    cv2.rectangle(overlap_mask, (x1, y1), (x2, y2), 255, -1)

            for item in sam_dict.get(page_idx, []) or []:
                bbox_px = item.get("bbox_px")
                if bbox_px and len(bbox_px) == 4:
                    x1, y1, x2, y2 = [int(v) for v in bbox_px]
                else:
                    bbox_n = item.get("bbox")
                    if not bbox_n or len(bbox_n) != 4:
                        continue
                    x1 = int(float(bbox_n[0]) * img_w)
                    y1 = int(float(bbox_n[1]) * img_h)
                    x2 = int(float(bbox_n[2]) * img_w)
                    y2 = int(float(bbox_n[3]) * img_h)

                pad = 3
                x1 = max(0, x1 - pad)
                y1 = max(0, y1 - pad)
                x2 = min(img_w, x2 + pad)
                y2 = min(img_h, y2 + pad)
                if x2 > x1 and y2 > y1:
                    cv2.rectangle(overlap_mask, (x1, y1), (x2, y2), 255, -1)

            text_mask_path = pinfo.get("text_mask_path")
            if text_mask_path and os.path.exists(text_mask_path):
                text_mask = cv2.imread(text_mask_path, cv2.IMREAD_GRAYSCALE)
                if text_mask is not None:
                    if text_mask.shape[:2] != (img_h, img_w):
                        text_mask = cv2.resize(text_mask, (img_w, img_h), interpolation=cv2.INTER_NEAREST)
                    overlap_mask = np.maximum(overlap_mask, text_mask)

            if not np.any(overlap_mask):
                return None

            overlap_mask = cv2.dilate(overlap_mask, np.ones((5, 5), dtype=np.uint8), iterations=1)
            cleaned_bg = cv2.inpaint(bg_img, overlap_mask, 3, cv2.INPAINT_TELEA)

            out_dir = base_dir / "clean_bg_overlay_free"
            out_dir.mkdir(parents=True, exist_ok=True)
            out_path = out_dir / f"bg_{page_idx+1:03d}.png"
            if cv2.imwrite(str(out_path), cleaned_bg):
                return str(out_path)
            return None

        async def _process_inpainting(pinfo):
            page_idx = pinfo.get("page_idx", 0)
            img_path = pinfo.get("path")
            no_text_path = pinfo.get("no_text_path")
            
            clean_bg_path = base_dir / "clean_bg" / f"bg_{page_idx+1:03d}.png"
            clean_bg_path.parent.mkdir(parents=True, exist_ok=True)

            clean_bg_lite_path = base_dir / "clean_bg_lite" / f"bg_{page_idx+1:03d}.png"
            clean_bg_lite_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 将结果路径回写到 pinfo，供后续步骤使用
            pinfo["clean_bg_path"] = str(clean_bg_path)
            pinfo["clean_bg_lite_path"] = str(clean_bg_lite_path)

            # 1) 基于 OCR 文本框生成 no_text 图 + text mask
            merged_no_text_path = None
            merged_mask_path = None
            img_w = img_h = 0
            if img_path and os.path.exists(img_path):
                try:
                    img_bgr = cv2.imread(img_path, cv2.IMREAD_COLOR)
                    if img_bgr is not None:
                        img_h, img_w = img_bgr.shape[:2]
                        mask_img = img_bgr.copy()
                        text_mask = np.zeros((img_h, img_w), dtype=np.uint8)

                        # OCR 文本框（已是 [y1, x1, y2, x2]）
                        merged_boxes: List[List[float]] = []
                        for it in pinfo.get("vlm_data", []) or []:
                            bbox_n = it.get("bbox")
                            if bbox_n and len(bbox_n) == 4:
                                merged_boxes.append(bbox_n)

                        if merged_boxes:
                            for bbox_n in merged_boxes:
                                y1n, x1n, y2n, x2n = bbox_n
                                x1 = int(x1n * img_w)
                                y1 = int(y1n * img_h)
                                x2 = int(x2n * img_w)
                                y2 = int(y2n * img_h)
                                box_h = max(1, y2 - y1)
                                pad = max(4, int(0.15 * box_h))
                                x1 = max(0, x1 - pad)
                                y1 = max(0, y1 - pad)
                                x2 = min(img_w, x2 + pad)
                                y2 = min(img_h, y2 + pad)
                                cv2.rectangle(mask_img, (x1, y1), (x2, y2), (255, 255, 255), -1)
                                cv2.rectangle(text_mask, (x1, y1), (x2, y2), 255, -1)

                            merged_no_text_path = base_dir / "vlm_debug" / f"page_{page_idx+1:03d}_no_text_merged.png"
                            merged_mask_path = base_dir / "vlm_debug" / f"page_{page_idx+1:03d}_text_mask.png"
                            merged_no_text_path.parent.mkdir(parents=True, exist_ok=True)
                            cv2.imwrite(str(merged_no_text_path), mask_img)
                            cv2.imwrite(str(merged_mask_path), text_mask)
                            no_text_path = str(merged_no_text_path)
                            pinfo["no_text_path"] = no_text_path
                            pinfo["text_mask_path"] = str(merged_mask_path)
                except Exception as e:
                    log.warning(f"[pdf2ppt_qwenvl][Inpainting] page#{page_idx+1} merged no_text failed: {e}")

            # 2) 自适应轻量去字，用于 SAM 输入
            lite_ok = False
            if img_path and os.path.exists(img_path):
                if img_w == 0 or img_h == 0:
                    try:
                        img_bgr = cv2.imread(img_path, cv2.IMREAD_COLOR)
                        if img_bgr is not None:
                            img_h, img_w = img_bgr.shape[:2]
                    except Exception:
                        pass
                text_boxes = []
                for it in pinfo.get("vlm_data", []) or []:
                    bbox_n = it.get("bbox")
                    if bbox_n and len(bbox_n) == 4:
                        text_boxes.append(bbox_n)
                expanded_boxes = []
                if text_boxes and img_w > 0 and img_h > 0:
                    for bbox_n in text_boxes:
                        y1n, x1n, y2n, x2n = bbox_n
                        box_h_px = max(1, int((y2n - y1n) * img_h))
                        pad_px = max(4, int(0.15 * box_h_px))
                        exp = _expand_bbox_norm(bbox_n, img_w, img_h, pad_px)
                        if exp:
                            expanded_boxes.append(exp)
                if expanded_boxes:
                    try:
                        lite_ok = _adaptive_fill_text_regions(
                            img_path=img_path,
                            text_boxes_norm=expanded_boxes,
                            output_path=str(clean_bg_lite_path),
                            skip_boxes_norm=None,
                            pad=0,
                        )
                    except Exception as e:
                        log.warning(f"[pdf2ppt_qwenvl][Inpainting] page#{page_idx+1} adaptive fill failed: {e}")
            
            mask_path = pinfo.get("text_mask_path")
            # AI 编辑输入图：优先使用 merged no_text 图，其次原图
            edit_image_path = None
            if no_text_path and os.path.exists(no_text_path):
                edit_image_path = no_text_path
            elif img_path and os.path.exists(img_path):
                edit_image_path = img_path

            if state.use_ai_edit and api_key and edit_image_path:
                ratio_str = "16:9"
                try:
                    with Image.open(edit_image_path) as tmp_img:
                        ratio_str = get_closest_aspect_ratio(tmp_img.width, tmp_img.height)
                except Exception: pass
                
                inpainting_prompt = "Fill the masked areas with matching background. Remove text."
                log.info(
                    f"[pdf2ppt_qwenvl][Inpainting] page#{page_idx+1} edit_image={edit_image_path}, "
                    f"mask={mask_path if (mask_path and os.path.exists(mask_path)) else 'None'}"
                )
                
                async with sem:
                    await _call_image_api_with_retry(
                        lambda: generate_or_edit_and_save_image_async(
                            prompt=inpainting_prompt,
                            save_path=str(clean_bg_path),
                            api_url=api_url,
                            api_key=api_key,
                            model=model_name,
                            use_edit=True,
                            image_path=edit_image_path,
                            mask_path=mask_path if (mask_path and os.path.exists(mask_path)) else None,
                            aspect_ratio=ratio_str,
                            resolution="2K"
                        )
                    )

            # 2) 降级逻辑：优先使用自适应轻量图，其次 no_text，最后原图
            final_bg_path = None
            if os.path.exists(clean_bg_path):
                final_bg_path = str(clean_bg_path)
            elif lite_ok and clean_bg_lite_path.exists():
                final_bg_path = str(clean_bg_lite_path)
            elif no_text_path and os.path.exists(no_text_path):
                final_bg_path = str(no_text_path)
            elif img_path and os.path.exists(img_path):
                final_bg_path = str(img_path)

            if final_bg_path:
                cleaned_bg_path = _cleanup_bg_overlay(page_idx, pinfo, final_bg_path, img_w, img_h)
                # 保持 clean_bg_path 指向“原始AI编辑背景图”；overlay_free 单独记录
                if cleaned_bg_path:
                    pinfo["clean_bg_overlay_free_path"] = cleaned_bg_path

        tasks = [_process_inpainting(p) for p in vlm_pages]
        if tasks:
            log.info(f"[pdf2ppt_qwenvl][Inpainting] starting {len(tasks)} tasks")
            await asyncio.gather(*tasks)
            
        return state

    # --- 主处理节点 ---
    async def parallel_processing_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        串行流程：
        1. VLM OCR
        2. SAM3 分割 + 图标去背
        3. AI 背景编辑
        """
        import time
        start_time = time.time()

        state = await vlm_recognition_node(state)
        state = await slides_segment_hint_node(state)
        state = await slides_sam_node(state)
        sam_pages = getattr(state, "sam_pages", [])
        state = await slides_layout_bg_remove_node(state, sam_pages=sam_pages)
        state = await slides_inpainting_node(state)
        
        log.info(f"[pdf2ppt_qwenvl] Processing finished in {time.time() - start_time:.2f}s")
        return state

    async def slides_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        生成 PPT：
        1. 整合 VLM, SAM3 结果
        2. 渲染页面 (Inpainting 已经在上一步并行完成)
        """
        vlm_pages = getattr(state, "vlm_pages", []) or []
        sam_pages = getattr(state, "sam_pages", []) or []

        if not vlm_pages:
            log.error("[pdf2ppt_qwenvl] no vlm_pages, abort PPT generation")
            return state

        # Indexing
        sam_dict = {p.get("page_idx", 0): p.get("layout_items", []) for p in sam_pages}

        prs = Presentation()
        prs.slide_width = Inches(ppt_tool.SLIDE_W_IN)
        prs.slide_height = Inches(ppt_tool.SLIDE_H_IN)
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height

        base_dir = Path(_ensure_result_path(state))

        # 辅助几何函数
        def _bbox_area(bbox): return max(0, bbox[2]-bbox[0]) * max(0, bbox[3]-bbox[1])
        def _get_intersection_area(b1, b2):
            x1,y1,x2,y2 = max(b1[0],b2[0]), max(b1[1],b2[1]), min(b1[2],b2[2]), min(b1[3],b2[3])
            return max(0, x2-x1) * max(0, y2-y1)
        def _is_inside(inner, outer, th=0.9):
            ia = _bbox_area(inner)
            return (ia > 0) and ((_get_intersection_area(inner, outer) / ia) >= th)
        
        for pinfo in vlm_pages:
            page_idx = pinfo.get("page_idx", 0)
            img_path = pinfo.get("path")
            vlm_data = pinfo.get("vlm_data", [])
            # 背景图优先级：clean_bg(模型编辑) > overlay_free > clean_bg_lite > no_text > 原图
            bg_candidates = [
                pinfo.get("clean_bg_path"),
                pinfo.get("clean_bg_overlay_free_path"),
                pinfo.get("clean_bg_lite_path"),
                pinfo.get("no_text_path"),
                img_path,
            ]
            clean_bg_path = next(
                (p for p in bg_candidates if p and os.path.exists(p)),
                None,
            )
            if clean_bg_path:
                log.info(
                    f"[pdf2ppt_qwenvl][PPT] page#{page_idx+1} background={clean_bg_path}"
                )

            if not img_path or not os.path.exists(img_path): continue

            try:
                pil_img = Image.open(img_path)
                w0, h0 = pil_img.size
            except Exception: continue

            scale_x = slide_w_emu / w0
            scale_y = slide_h_emu / h0

            # 2. VLM Text Filtering
            final_text_lines = []
            for it in vlm_data:
                # it: {'bbox': [y1n, x1n, y2n, x2n], 'text': ...}  (0-1 norm)
                bbox_n = it.get("bbox")
                if not bbox_n: continue
                y1n, x1n, y2n, x2n = bbox_n
                x1, y1, x2, y2 = int(x1n*w0), int(y1n*h0), int(x2n*w0), int(y2n*h0)
                l_bbox = [x1, y1, x2, y2]
                
                raw_pt_obj = ppt_tool.estimate_font_pt(l_bbox, img_h_px=h0, body_h_px=None)
                raw_pt = raw_pt_obj.pt if hasattr(raw_pt_obj, "pt") else raw_pt_obj

                l_type = "body"
                if raw_pt > 18:
                    l_type = "title"

                final_text_lines.append((l_bbox, it.get("text", ""), 1.0, l_type, raw_pt))

            # 3. SAM3 Icons Filtering
            raw_sam = sam_dict.get(page_idx, [])
            final_sam = []
            for item in raw_sam:
                s_bbox = item.get("bbox_px")
                if not s_bbox:
                    continue
                if any(_is_inside(line[0], s_bbox) for line in final_text_lines):
                    continue
                
                final_sam.append(item)

            # 渲染 PPT 页面
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Background
            if clean_bg_path and os.path.exists(clean_bg_path):
                try: slide.shapes.add_picture(clean_bg_path, 0, 0, prs.slide_width, prs.slide_height)
                except: pass
            
            # SAM Icons
            for s in final_sam:
                path = s.get("fg_png_path") or s.get("png_path")
                if path and os.path.exists(path):
                    bx = s["bbox_px"]
                    slide.shapes.add_picture(path,
                        ppt_tool.px_to_emu(bx[0], scale_x), ppt_tool.px_to_emu(bx[1], scale_y),
                        ppt_tool.px_to_emu(bx[2]-bx[0], scale_x), ppt_tool.px_to_emu(bx[3]-bx[1], scale_y))

            # Text
            for line in final_text_lines:
                bbox, text, _, l_type, raw_pt = line
                left = ppt_tool.px_to_emu(bbox[0], scale_x)
                top = ppt_tool.px_to_emu(bbox[1], scale_y)
                w = ppt_tool.px_to_emu(bbox[2]-bbox[0], scale_x)
                h = ppt_tool.px_to_emu(bbox[3]-bbox[1], scale_y)
                
                tb = slide.shapes.add_textbox(left, top, w, h)
                p = tb.text_frame.paragraphs[0]
                p.text = text
                p.font.size = Pt(raw_pt if raw_pt > 5 else 12)
                p.font.bold = (l_type == "title")
                p.font.color.rgb = RGBColor(0,0,0)

        out_path = base_dir / "pdf2ppt_qwenvl_output.pptx"
        prs.save(str(out_path))
        state.ppt_path = str(out_path)
        log.info(f"[pdf2ppt_qwenvl] PPT Generated: {out_path}")

        return state

    nodes = {
        "_start_": _init_result_path,
        "pdf_to_images": pdf_to_images_node,
        "parallel_processing": parallel_processing_node,
        "slides_ppt_generation": slides_ppt_generation_node,
        "_end_": lambda s: s,
    }

    edges = [
        ("pdf_to_images", "parallel_processing"),
        ("parallel_processing", "slides_ppt_generation"),
        ("slides_ppt_generation", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges)
    builder.add_edge("_start_", "pdf_to_images")
    return builder
